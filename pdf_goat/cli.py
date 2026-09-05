"""Local PDF editing and inspection CLI.

It writes JSON when piped or with ``--agent`` and plain text on a TTY.
The CLI records each parsed command run in SQLite, except ``capabilities`` and ``jobs``.
"""

import argparse
import json
import os
import re
import sys
import time
import unicodedata
from itertools import islice
from pathlib import Path

from .layout import extract_page_layout

HOME = Path(os.environ.get("PDF_GOAT_HOME", Path.home() / ".pdf-goat"))
DB_PATH = HOME / "ledger.db"

_DEFAULT_PAGE_WINDOW = 25

# argparse passes every help string through gettext, and each lookup probes
# the locale directory for a catalog that neither Python nor this CLI ships:
# 339 lookups, 1,356 stat calls, and 5 ms of the 8 ms parser build per run.
argparse._ = lambda message: message
argparse.ngettext = lambda singular, plural, n: singular if n == 1 else plural


# --------------------------------------------------------------------------- #
# Ledger
# --------------------------------------------------------------------------- #
def _db():
    import sqlite3

    HOME.mkdir(parents=True, exist_ok=True)
    conn = sqlite3.connect(DB_PATH)
    conn.execute(
        """CREATE TABLE IF NOT EXISTS jobs(
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            ts TEXT NOT NULL,
            verb TEXT NOT NULL,
            status TEXT NOT NULL,
            inputs TEXT,
            outputs TEXT,
            detail TEXT,
            message TEXT,
            duration_ms INTEGER
        )"""
    )
    return conn


def record_job(verb, status, inputs, outputs, detail, message, duration_ms):
    with _db() as conn:
        cur = conn.execute(
            "INSERT INTO jobs(ts,verb,status,inputs,outputs,detail,message,duration_ms)"
            " VALUES(?,?,?,?,?,?,?,?)",
            (
                time.strftime("%Y-%m-%dT%H:%M:%S"),
                verb,
                status,
                json.dumps(inputs),
                json.dumps(outputs),
                json.dumps(detail, default=str),
                message,
                duration_ms,
            ),
        )
        return cur.lastrowid


# --------------------------------------------------------------------------- #
# Helpers
# --------------------------------------------------------------------------- #
class PdfGoatError(Exception):
    pass


class PdfGoatArgumentParser(argparse.ArgumentParser):
    def error(self, message):
        if message.startswith("unrecognized arguments:"):
            safe_message = "unrecognized arguments"
        else:
            head, separator, detail = message.partition(": ")
            safe_detail = detail.startswith(("not allowed with argument ", "expected "))
            safe_message = (
                message
                if not separator
                or safe_detail
                or head == "the following arguments are required"
                else head
            )
        raise PdfGoatError(safe_message)


def resolve(path):
    p = Path(path).expanduser()
    if not p.exists():
        raise PdfGoatError(f"file not found: {path}")
    return p.resolve()


def default_out(src, suffix, ext="pdf"):
    src = Path(src)
    return str(src.with_name(f"{src.stem}.{suffix}.{ext}"))


def ensure_parent(path):
    Path(path).expanduser().resolve().parent.mkdir(parents=True, exist_ok=True)
    return str(Path(path).expanduser().resolve())


class AtomicOutput:
    """Write through a sibling `.part` file so a failed write leaves nothing.

    For verbs that write a whole file in one step. An in-place `.save(out)` on
    an open pikepdf document keeps pikepdf's own same-file guard instead.
    """

    def __init__(self, path):
        self.path = str(path)
        self.partial = f"{self.path}.part"

    def __enter__(self):
        return self.partial

    def __exit__(self, error_type, error, traceback):
        try:
            if error_type is None:
                os.replace(self.partial, self.path)
        finally:
            Path(self.partial).unlink(missing_ok=True)


def _save_pdf(doc, out, **options):
    """Write a MuPDF document without the duplicate-object scan.

    ``garbage=3`` and ``4`` compare every object pair, which costs 5 to 8
    seconds on a 1,000-page file for a 1% size gain. Object streams keep the
    output as small as a modern source.
    """
    doc.save(out, garbage=2, deflate=True, use_objstms=1, **options)


def parse_pages(spec, n):
    """Parse '2-5,9' (1-based, inclusive) into ordered 0-based indices."""
    out = []
    for part in spec.split(","):
        part = part.strip()
        if not part:
            continue
        if "-" in part:
            a, b = part.split("-", 1)
            a, b = int(a), int(b)
            step = 1 if b >= a else -1
            out.extend(range(a, b + step, step))
        else:
            out.append(int(part))
    for p in out:
        if p < 1 or p > n:
            raise PdfGoatError(f"page {p} is outside the range 1 to {n}")
    return [p - 1 for p in out]


def human_size(num):
    for unit in ("B", "KB", "MB", "GB"):
        if abs(num) < 1024:
            return f"{num:.1f}{unit}" if unit != "B" else f"{num}B"
        num /= 1024
    return f"{num:.1f}TB"


def weasyprint_bin():
    import shutil

    wp = shutil.which("weasyprint")
    if not wp:
        raise PdfGoatError("from-html and from-md require weasyprint on PATH")
    return wp


def run_weasyprint(args, stdin=None):
    import subprocess

    proc = subprocess.run(
        [weasyprint_bin(), *args],
        input=stdin,
        text=True,
        capture_output=True,
        check=False,
    )
    if proc.returncode != 0:
        raise PdfGoatError(
            f"weasyprint failed: {proc.stderr.strip() or proc.returncode}"
        )


DEFAULT_CSS = """
@page { size: A4; margin: 2cm; }
body { font-family: -apple-system, "Helvetica Neue", Arial, sans-serif;
  font-size: 11pt; line-height: 1.55; color: #1a1a1a; }
h1, h2, h3, h4 { font-weight: 600; line-height: 1.25; margin: 1.4em 0 .5em; color: #111; }
h1 { font-size: 1.9em; } h2 { font-size: 1.45em; } h3 { font-size: 1.2em; }
p { margin: 0 0 .8em; }
a { color: #0b5cad; text-decoration: none; }
ul, ol { margin: 0 0 .8em 1.4em; }
li { margin: .2em 0; }
code { font-family: "SF Mono", ui-monospace, Menlo, monospace; font-size: .9em;
  background: #f3f4f6; padding: .1em .3em; border-radius: 3px; }
pre { background: #f3f4f6; padding: .9em 1em; border-radius: 6px; overflow: auto; }
pre code { background: none; padding: 0; }
blockquote { margin: 0 0 .8em; padding: .2em 1em; border-left: 3px solid #d0d5dd; color: #475467; }
table { border-collapse: collapse; width: 100%; margin: 0 0 1em; font-size: .95em; }
th, td { border: 1px solid #d0d5dd; padding: .45em .7em; text-align: left; }
th { background: #f9fafb; font-weight: 600; }
img { max-width: 100%; }
hr { border: none; border-top: 1px solid #e4e7ec; margin: 1.5em 0; }
"""


# --------------------------------------------------------------------------- #
# Verbs
# --------------------------------------------------------------------------- #
def cmd_info(a):
    import pymupdf

    src = resolve(a.file)
    doc = pymupdf.open(src)
    has_text = False
    sizes = []
    for page in doc:
        sizes.append((round(page.rect.width, 1), round(page.rect.height, 1)))
        if not has_text and page.get_text("text").strip():
            has_text = True
    uniq = sorted(set(sizes))
    result = {
        "verb": "info",
        "inputs": [str(src)],
        "outputs": [],
        "pages": doc.page_count,
        "encrypted": bool(doc.is_encrypted),
        "needs_password": bool(doc.needs_pass),
        "has_forms": bool(doc.is_form_pdf),
        "form_field_count": int(doc.is_form_pdf) if doc.is_form_pdf else 0,
        "has_text": has_text,
        "page_sizes_pt": [{"width": w, "height": h} for w, h in uniq],
        "file_size_bytes": src.stat().st_size,
        "metadata": {k: v for k, v in doc.metadata.items() if v},
    }
    doc.close()
    return result


def _subcommand_parsers(parser):
    for action in parser._actions:
        if isinstance(action, argparse._SubParsersAction):
            return action.choices
    return {}


def _argument_schema(action):
    value_type = {int: "integer", float: "number"}.get(action.type, "string")
    if action.nargs == 0:
        value_type = "boolean"
    item = {
        "name": action.dest,
        "flags": action.option_strings,
        "required": action.required,
        "type": value_type,
    }
    if isinstance(action, argparse._AppendAction):
        item["repeatable"] = True
    if action.nargs is not None:
        item["nargs"] = action.nargs
    if action.choices is not None:
        item["choices"] = list(action.choices)
    if action.default is not None and action.default is not argparse.SUPPRESS:
        item["default"] = action.default
    if action.help:
        item["help"] = action.help
    return item


def _parser_arguments(parser):
    return [
        _argument_schema(action)
        for action in parser._actions
        if not isinstance(action, (argparse._HelpAction, argparse._SubParsersAction))
    ]


def _command_schema(parser):
    groups = [
        {
            "required": group.required,
            "arguments": [action.dest for action in group._group_actions],
        }
        for group in parser._mutually_exclusive_groups
        if group._group_actions
    ]
    schema = {
        "command": parser.prog,
        "arguments": _parser_arguments(parser),
        "commands": {
            name: _command_schema(child)
            for name, child in _subcommand_parsers(parser).items()
        },
    }
    if groups:
        schema["mutually_exclusive_groups"] = groups
    return schema


def _leaf_command_count(parser):
    children = _subcommand_parsers(parser)
    if not children:
        return 1
    return sum(_leaf_command_count(child) for child in children.values())


def cmd_capabilities(a):
    parser = a.root_parser
    top_level = _subcommand_parsers(parser)
    if a.family and a.family not in top_level:
        raise PdfGoatError(f"unknown top-level command: {a.family}")
    families = sorted(
        name for name, command in top_level.items() if _subcommand_parsers(command)
    )
    commands = sorted(set(top_level) - set(families))
    schemas = {a.family: _command_schema(top_level[a.family])} if a.family else {}
    return {
        "verb": "capabilities",
        "inputs": [],
        "outputs": [],
        "schema_version": 1,
        "mode": "standalone",
        "agent_json": True,
        "root_arguments": _parser_arguments(parser),
        "families": families,
        "commands": commands,
        "requested_command": a.family,
        "command_count": _leaf_command_count(parser),
        "schemas": schemas,
    }


def _page_text_and_words(page):
    """Extract both forms from one text page; MuPDF builds it once per page."""
    import pymupdf

    textpage = page.get_textpage(flags=pymupdf.TEXTFLAGS_TEXT)
    return (
        page.get_text("text", textpage=textpage),
        page.get_text("words", textpage=textpage),
    )


def _page_inventory(page, index):
    text, words = _page_text_and_words(page)
    return {
        "page": index + 1,
        "label": page.get_label(),
        "width_pt": round(page.rect.width, 1),
        "height_pt": round(page.rect.height, 1),
        "rotation": page.rotation,
        "text_chars": len(text),
        "word_count": len(words),
        "image_count": len(page.get_images(full=True)),
        "link_count": len(page.get_links()),
        "annotation_count": sum(1 for _ in (page.annots() or ())),
        "form_field_count": sum(1 for _ in (page.widgets() or ())),
    }


def cmd_inspect(a):
    import pymupdf

    src = resolve(a.file)
    doc = pymupdf.open(src)
    if doc.needs_pass:
        doc.close()
        raise PdfGoatError("inspect cannot read an encrypted PDF without a password")
    page_count = doc.page_count
    if a.start_page < 1 or a.start_page > page_count:
        doc.close()
        raise PdfGoatError(f"--start-page must be between 1 and {page_count}")
    if a.limit < 1 or a.limit > 100:
        doc.close()
        raise PdfGoatError("--limit must be between 1 and 100")

    start = a.start_page - 1
    end = min(start + a.limit, doc.page_count)
    pages = [_page_inventory(doc[index], index) for index in range(start, end)]
    total_pages = doc.page_count
    doc.close()
    return {
        "verb": "inspect",
        "inputs": [str(src)],
        "outputs": [],
        "total_pages": total_pages,
        "start_page": a.start_page,
        "pages": pages,
        "next_page": end + 1 if end < total_pages else None,
        "truncated": end < total_pages,
    }


def _external_link_kinds():
    import pymupdf

    return pymupdf.LINK_URI, pymupdf.LINK_LAUNCH, pymupdf.LINK_GOTOR


def _link_preflight(page):
    import pymupdf

    external_links = 0
    unsafe_links = 0
    external_kinds = _external_link_kinds()
    for link in page.get_links():
        if link["kind"] not in external_kinds:
            continue
        external_links += 1
        uri = link.get("uri") or link.get("file") or ""
        scheme = uri.partition(":")[0].lower()
        if link["kind"] != pymupdf.LINK_URI or scheme not in {
            "http",
            "https",
            "mailto",
        }:
            unsafe_links += 1
    return external_links, unsafe_links


def _page_preflight(page):
    external_links, unsafe_links = _link_preflight(page)
    return {
        "annotations": sum(1 for _ in (page.annots() or ())),
        "form_fields": sum(1 for _ in (page.widgets() or ())),
        "external_links": external_links,
        "unsafe_links": unsafe_links,
        "empty": not page.get_fonts() and not page.get_images(full=True),
    }


def _document_preflight(doc):
    profile = {
        "annotations": 0,
        "form_fields": 0,
        "external_links": 0,
        "unsafe_links": 0,
        "empty_pages": [],
    }
    for index, page in enumerate(doc):
        page_profile = _page_preflight(page)
        for key in ("annotations", "form_fields", "external_links", "unsafe_links"):
            profile[key] += page_profile[key]
        if page_profile["empty"]:
            profile["empty_pages"].append(index + 1)
    return profile


def _file_attachment_annotations(doc):
    import pymupdf

    for page_index in range(doc.page_count):
        page = doc[page_index]
        for annotation in page.annots() or ():
            if annotation.type[0] == pymupdf.PDF_ANNOT_FILE_ATTACHMENT:
                yield page, annotation, annotation.file_info.get("filename") or ""


def _active_content_count(structure):
    return (
        structure["root_actions"]
        + structure["javascript_actions"]
        + structure["launch_actions"]
    )


def _direct_pdf_dictionaries(obj):
    import pikepdf

    stack = [obj]
    while stack:
        item = stack.pop()
        match item:
            case pikepdf.Dictionary():
                yield item
                children = item.values()
            case pikepdf.Array():
                children = item
            case _:
                continue
        stack.extend(
            child
            for child in children
            if isinstance(child, (pikepdf.Dictionary, pikepdf.Array))
            and not child.is_indirect
        )


def _pdf_object_actions(obj):
    root_actions = javascript_actions = launch_actions = signatures = 0
    for dictionary in _direct_pdf_dictionaries(obj):
        action = str(dictionary["/S"]) if "/S" in dictionary else ""
        root_actions += int("/OpenAction" in dictionary) + int("/AA" in dictionary)
        javascript_actions += int("/JS" in dictionary or action == "/JavaScript")
        launch_actions += int(action == "/Launch")
        signatures += int("/FT" in dictionary and str(dictionary["/FT"]) == "/Sig")
    return root_actions, javascript_actions, launch_actions, signatures


def _structure_preflight(src):
    import pikepdf

    profile = {
        "root_actions": 0,
        "javascript_actions": 0,
        "launch_actions": 0,
        "signatures": 0,
        "xfa": False,
        "tagged": False,
        "language": None,
    }
    with pikepdf.open(src) as pdf:
        root = pdf.Root
        profile["tagged"] = "/StructTreeRoot" in root
        if "/MarkInfo" in root and "/Marked" in root.MarkInfo:
            profile["tagged"] = profile["tagged"] or bool(root.MarkInfo.Marked)
        profile["language"] = str(root.Lang) if "/Lang" in root else None
        profile["xfa"] = "/AcroForm" in root and "/XFA" in root.AcroForm
        for obj in pdf.objects:
            root_actions, javascript_actions, launch_actions, signatures = (
                _pdf_object_actions(obj)
            )
            profile["root_actions"] += root_actions
            profile["javascript_actions"] += javascript_actions
            profile["launch_actions"] += launch_actions
            profile["signatures"] += signatures
    return profile


def _preflight_findings(content, structure, title, attachment_count):
    findings = []
    active_count = _active_content_count(structure)
    if active_count:
        findings.append(
            {
                "code": "active_content",
                "severity": "warning",
                "message": "The PDF contains document-open actions, JavaScript actions, or launch actions.",
                "count": active_count,
            }
        )
    if content["unsafe_links"]:
        findings.append(
            {
                "code": "unsafe_links",
                "severity": "warning",
                "message": "The PDF contains links with schemes other than HTTP, HTTPS, or mailto.",
                "count": content["unsafe_links"],
            }
        )
    if attachment_count:
        findings.append(
            {
                "code": "attachments",
                "severity": "info",
                "message": "The PDF contains embedded or attached files.",
                "count": attachment_count,
            }
        )
    if structure["xfa"]:
        findings.append(
            {
                "code": "xfa",
                "severity": "warning",
                "message": "The PDF contains an unsupported XFA form.",
            }
        )
    if not structure["tagged"]:
        findings.append(
            {
                "code": "untagged",
                "severity": "info",
                "message": "The PDF has no tag tree.",
            }
        )
    if not title:
        findings.append(
            {
                "code": "missing_title",
                "severity": "info",
                "message": "The PDF has no title metadata.",
            }
        )
    if not structure["language"]:
        findings.append(
            {
                "code": "missing_language",
                "severity": "info",
                "message": "The PDF has no document language.",
            }
        )
    if content["empty_pages"]:
        findings.append(
            {
                "code": "empty_pages",
                "severity": "info",
                "message": "Some pages declare no font and no image.",
                "pages": content["empty_pages"][:100],
                "count": len(content["empty_pages"]),
            }
        )
    return findings


def cmd_preflight(a):
    import pymupdf

    src = resolve(a.file)
    doc = pymupdf.open(src)
    encrypted = bool(doc.is_encrypted)
    if doc.needs_pass:
        doc.close()
        return {
            "verb": "preflight",
            "inputs": [str(src)],
            "outputs": [],
            "risk": "unknown",
            "encrypted": encrypted,
            "needs_password": True,
            "findings": [
                {
                    "code": "encrypted",
                    "severity": "info",
                    "message": "A password is required before content checks can run.",
                }
            ],
        }

    content = _document_preflight(doc)
    attachment_count = len(doc.embfile_names()) + sum(
        1 for _ in _file_attachment_annotations(doc)
    )
    title = doc.metadata.get("title")
    pages = doc.page_count
    doc.close()
    structure = _structure_preflight(src)
    findings = _preflight_findings(content, structure, title, attachment_count)
    active_count = _active_content_count(structure)
    risk = (
        "high"
        if active_count or content["unsafe_links"]
        else "medium"
        if attachment_count or structure["xfa"]
        else "low"
    )
    return {
        "verb": "preflight",
        "inputs": [str(src)],
        "outputs": [],
        "risk": risk,
        "encrypted": encrypted,
        "needs_password": False,
        "pages": pages,
        "annotations": content["annotations"],
        "form_fields": content["form_fields"],
        "signatures": structure["signatures"],
        "external_links": content["external_links"],
        "unsafe_links": content["unsafe_links"],
        "attachments": attachment_count,
        "javascript_actions": structure["javascript_actions"],
        "launch_actions": structure["launch_actions"],
        "findings": findings,
    }


def cmd_merge(a):
    import pikepdf

    srcs = [resolve(f) for f in a.files]
    out = ensure_parent(a.output or "merged.pdf")
    merged = pikepdf.Pdf.new()
    total = 0
    for s in srcs:
        with pikepdf.open(s) as pdf:
            merged.pages.extend(pdf.pages)
            total += len(pdf.pages)
    merged.save(out)
    merged.close()
    return {
        "verb": "merge",
        "inputs": [str(s) for s in srcs],
        "outputs": [out],
        "merged_pages": total,
    }


def cmd_split(a):
    import pikepdf

    src = resolve(a.file)
    every = max(1, a.every)
    outdir = Path(a.outdir or f"{src.stem}_split").expanduser().resolve()
    outdir.mkdir(parents=True, exist_ok=True)
    outputs = []
    with pikepdf.open(src) as pdf:
        n = len(pdf.pages)
        for idx, start in enumerate(range(0, n, every)):
            part = pikepdf.Pdf.new()
            for p in range(start, min(start + every, n)):
                part.pages.append(pdf.pages[p])
            out = outdir / f"{src.stem}_{idx + 1:03d}.pdf"
            with AtomicOutput(out) as partial:
                part.save(partial)
            part.close()
            outputs.append(str(out))
    return {
        "verb": "split",
        "inputs": [str(src)],
        "outputs": outputs,
        "every": every,
        "parts": len(outputs),
    }


def cmd_extract(a):
    import pikepdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "extract"))
    with pikepdf.open(src) as pdf:
        idxs = parse_pages(a.pages, len(pdf.pages))
        new = pikepdf.Pdf.new()
        for i in idxs:
            new.pages.append(pdf.pages[i])
        with AtomicOutput(out) as partial:
            new.save(partial)
        new.close()
    return {
        "verb": "extract",
        "inputs": [str(src)],
        "outputs": [out],
        "pages": [i + 1 for i in idxs],
    }


def cmd_delete(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "deleted"))
    with pymupdf.open(src) as doc:
        n = doc.page_count
        drop = sorted(set(parse_pages(a.pages, n)))
        if len(drop) == n:
            raise PdfGoatError("--pages would remove every page")
        if drop:  # MuPDF prints "nothing to delete" to stdout on an empty list
            doc.delete_pages(drop)
        with AtomicOutput(out) as partial:
            _save_pdf(doc, partial)
    return {
        "verb": "delete",
        "inputs": [str(src)],
        "outputs": [out],
        "deleted_pages": [i + 1 for i in drop],
        "remaining_pages": n - len(drop),
    }


def cmd_reorder(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "reordered"))
    with pymupdf.open(src) as doc:
        n = doc.page_count
        idxs = parse_pages(a.order, n)
        if sorted(idxs) != list(range(n)):
            raise PdfGoatError(
                f"--order must list each of the {n} pages exactly once; "
                f"received {len(idxs)} entries"
            )
        doc.select(idxs)
        with AtomicOutput(out) as partial:
            _save_pdf(doc, partial)
    return {
        "verb": "reorder",
        "inputs": [str(src)],
        "outputs": [out],
        "order": [i + 1 for i in idxs],
    }


def cmd_rotate(a):
    import pikepdf

    if a.deg % 90 != 0:
        raise PdfGoatError("--deg must be a multiple of 90")
    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "rotated"))
    with pikepdf.open(src) as pdf:
        n = len(pdf.pages)
        idxs = set(parse_pages(a.pages, n)) if a.pages else set(range(n))
        for i in idxs:
            pdf.pages[i].rotate(a.deg, relative=True)
        pdf.save(out)
    return {
        "verb": "rotate",
        "inputs": [str(src)],
        "outputs": [out],
        "rotated_pages": sorted(i + 1 for i in idxs),
        "deg": a.deg,
    }


def cmd_render(a):
    import pymupdf

    src = resolve(a.file)
    outdir = Path(a.outdir or f"{src.stem}_render").expanduser().resolve()
    clip = pymupdf.Rect(parse_rect(a.clip)) if a.clip else None
    if clip is not None:
        clip.normalize()
        if clip.is_empty:
            raise PdfGoatError("--clip must have positive area")
    with pymupdf.open(src) as doc:
        indices = list(
            parse_pages(a.pages, doc.page_count) if a.pages else range(doc.page_count)
        )
        page_clips = []
        for index in indices:
            page_clip = doc[index].rect & clip if clip is not None else None
            if page_clip is not None and page_clip.is_empty:
                raise PdfGoatError(f"--clip does not overlap page {index + 1}")
            page_clips.append(page_clip)

        outdir.mkdir(parents=True, exist_ok=True)
        outputs = []
        for index, page_clip in zip(indices, page_clips, strict=True):
            pixmap = doc[index].get_pixmap(dpi=a.dpi, clip=page_clip)
            out = outdir / f"{src.stem}_p{index + 1:03d}.{a.format}"
            pixmap.save(out)
            outputs.append(str(out))
    return {
        "verb": "render",
        "inputs": [str(src)],
        "outputs": outputs,
        "dpi": a.dpi,
        "format": a.format,
        "clip": list(clip) if clip is not None else None,
    }


def _normalize_image(path, tmpdir):
    """Flatten RGBA, LA, and P images onto white because img2pdf does not accept alpha channels."""
    from PIL import Image

    img = Image.open(path)
    if img.mode in ("RGBA", "LA", "P"):
        rgba = img.convert("RGBA")
        bg = Image.new("RGB", rgba.size, (255, 255, 255))
        bg.paste(rgba, mask=rgba.split()[-1])
        out = Path(tmpdir) / (Path(path).stem + ".norm.png")
        bg.save(out)
        return str(out)
    return str(path)


def cmd_from_images(a):
    import tempfile

    import img2pdf

    srcs = [resolve(f) for f in a.images]
    out = ensure_parent(a.output or "images.pdf")
    with tempfile.TemporaryDirectory() as td:
        paths = [_normalize_image(s, td) for s in srcs]
        with open(out, "wb") as fh:
            fh.write(img2pdf.convert(paths))
    return {
        "verb": "from-images",
        "inputs": [str(s) for s in srcs],
        "outputs": [out],
        "image_count": len(srcs),
    }


def cmd_form_fields(a):
    from pypdf import PdfReader

    src = resolve(a.file)
    reader = PdfReader(src)
    fields = reader.get_fields() or {}
    out = []
    for name, f in fields.items():
        out.append(
            {
                "name": name,
                "type": f.get("/FT"),
                "value": f.get("/V"),
                "default": f.get("/DV"),
            }
        )
    return {
        "verb": "form-list",
        "inputs": [str(src)],
        "outputs": [],
        "field_count": len(out),
        "fields": out,
    }


def fill_pdf_form(src, data, out, flatten):
    import shutil
    import subprocess

    from pypdf import PdfReader, PdfWriter

    reader = PdfReader(src)
    writer = PdfWriter()
    writer.append(reader)
    for page in writer.pages:
        writer.update_page_form_field_values(page, data, auto_regenerate=False)
    with open(out, "wb") as fh:
        writer.write(fh)
    if flatten:
        qpdf = shutil.which("qpdf")
        if not qpdf:
            raise PdfGoatError("--flatten requires qpdf on PATH")
        tmp = Path(out).with_suffix(".flat.pdf")
        subprocess.run(
            [
                qpdf,
                "--generate-appearances",
                "--flatten-annotations=all",
                out,
                str(tmp),
            ],
            check=True,
            capture_output=True,
        )
        os.replace(tmp, out)
    return flatten


def cmd_form_fill(a):
    src = resolve(a.file)
    data = json.loads(Path(a.data).expanduser().read_text())
    out = ensure_parent(a.output or default_out(src, "filled"))
    flattened = fill_pdf_form(src, data, out, a.flatten)
    return {
        "verb": "form-fill",
        "inputs": [str(src)],
        "outputs": [out],
        "fields_set": list(data.keys()),
        "flattened": flattened,
    }


def cmd_redact(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "redacted"))
    pattern = re.compile(a.find)
    doc = pymupdf.open(src)
    hits = 0
    for page in doc:
        marked = False
        for x0, y0, x1, y1, word, *_ in page.get_text("words"):
            if pattern.search(word):
                page.add_redact_annot(pymupdf.Rect(x0, y0, x1, y1), fill=(0, 0, 0))
                hits += 1
                marked = True
        if marked:
            page.apply_redactions()
    doc.save(out, garbage=4, deflate=True)
    doc.close()
    return {
        "verb": "redact",
        "inputs": [str(src)],
        "outputs": [out],
        "pattern": a.find,
        "redactions": hits,
    }


def cmd_watermark(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "watermarked"))
    doc = pymupdf.open(src)
    for page in doc:
        w, h = page.rect.width, page.rect.height
        tl = pymupdf.get_text_length(a.text, fontsize=a.size)
        point = pymupdf.Point((w - tl) / 2, h / 2)
        tw = pymupdf.TextWriter(page.rect, color=(0.5, 0.5, 0.5))
        tw.append(point, a.text, fontsize=a.size)
        tw.write_text(page, opacity=a.opacity, morph=(point, pymupdf.Matrix(a.angle)))
    doc.save(out, garbage=4, deflate=True)
    doc.close()
    return {
        "verb": "watermark",
        "inputs": [str(src)],
        "outputs": [out],
        "text": a.text,
    }


def cmd_compress(a):
    import shutil
    import subprocess

    import pikepdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "compressed"))
    orig = src.stat().st_size
    gs = shutil.which("gs")
    source = src
    if gs:
        import tempfile

        handle, name = tempfile.mkstemp(prefix=f"{src.stem}.", suffix=".gs.pdf")
        os.close(handle)
        tmp = Path(name)
        proc = subprocess.run(
            [
                gs,
                "-sDEVICE=pdfwrite",
                f"-dPDFSETTINGS={a.level}",
                "-dCompatibilityLevel=1.5",
                "-dNOPAUSE",
                "-dBATCH",
                "-dQUIET",
                f"-sOutputFile={tmp}",
                str(src),
            ],
            capture_output=True,
            check=False,
        )
        if proc.returncode == 0 and tmp.stat().st_size > 0:
            source = tmp
        else:
            tmp.unlink(missing_ok=True)
    try:
        with pikepdf.open(source) as pdf, AtomicOutput(out) as partial:
            pdf.save(
                partial,
                linearize=True,
                compress_streams=True,
                object_stream_mode=pikepdf.ObjectStreamMode.generate,
            )
            new = Path(partial).stat().st_size
            # Keep the input bytes if recompression produces a larger file.
            kept_original = new >= orig
            if kept_original:
                shutil.copyfile(src, partial)
                new = orig
    finally:
        if source != src:
            Path(source).unlink(missing_ok=True)
    return {
        "verb": "compress",
        "inputs": [str(src)],
        "outputs": [out],
        "original_bytes": orig,
        "compressed_bytes": new,
        "ratio": round(new / orig, 3) if orig else None,
        "saved_bytes": orig - new,
        "linearized": not kept_original,
        "kept_original": kept_original,
        "used_ghostscript": bool(gs),
    }


def cmd_text(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output) if a.output else None
    doc = pymupdf.open(src)
    result = {
        "verb": "text",
        "inputs": [str(src)],
        "outputs": [out] if out else [],
    }
    if out and not a.layout:
        # The file is the output, so stream to it and report a page count
        # instead of holding the whole corpus twice more.
        chars = 0
        with Path(out).open("w") as handle:
            for index in range(doc.page_count):
                text = ("\n" if index else "") + doc[index].get_text("text")
                handle.write(text)
                chars += len(text)
        result["char_count"] = chars
        result["page_count"] = doc.page_count
        doc.close()
        return result
    pages = []
    for index in range(doc.page_count):
        page = doc[index]
        if a.layout:
            pages.append({"page": index + 1, **extract_page_layout(page)})
        else:
            pages.append({"page": index + 1, "text": page.get_text("text")})
    full = "\n".join(page["text"] for page in pages)
    doc.close()
    result["char_count"] = len(full)
    result["pages"] = pages
    if a.layout:
        result["mode"] = "layout"
    if out:  # layout mode only; plain text with -o returned above
        Path(out).write_text(full)
    return result


def cmd_from_html(a):
    src = resolve(a.file)
    out = ensure_parent(a.output or str(src.with_suffix(".pdf")))
    run_weasyprint([str(src), out, "-e", "utf-8"])
    return {
        "verb": "from-html",
        "inputs": [str(src)],
        "outputs": [out],
        "output_bytes": Path(out).stat().st_size,
    }


def cmd_from_md(a):
    import markdown

    src = resolve(a.file)
    out = ensure_parent(a.output or str(src.with_suffix(".pdf")))
    body = markdown.markdown(
        src.read_text(encoding="utf-8"),
        extensions=["tables", "fenced_code", "sane_lists", "toc"],
    )
    css = Path(a.css).expanduser().read_text() if a.css else DEFAULT_CSS
    html = (
        "<!DOCTYPE html><html><head><meta charset='utf-8'>"
        f"<title>{src.stem}</title><style>{css}</style></head><body>"
        f"{body}</body></html>"
    )
    run_weasyprint(
        ["-", out, "-e", "utf-8", "-u", src.parent.as_uri() + "/"], stdin=html
    )
    return {
        "verb": "from-md",
        "inputs": [str(src)],
        "outputs": [out],
        "output_bytes": Path(out).stat().st_size,
    }


def cmd_jobs(a):
    import sqlite3

    with _db() as conn:
        conn.row_factory = sqlite3.Row
        rows = conn.execute(
            "SELECT * FROM jobs ORDER BY id DESC LIMIT ?", (a.limit,)
        ).fetchall()
    jobs = []
    for r in rows:
        jobs.append(
            {
                "id": r["id"],
                "ts": r["ts"],
                "verb": r["verb"],
                "status": r["status"],
                "inputs": json.loads(r["inputs"]),
                "outputs": json.loads(r["outputs"]),
                "duration_ms": r["duration_ms"],
                "message": r["message"],
            }
        )
    return {
        "verb": "jobs",
        "inputs": [],
        "outputs": [],
        "count": len(jobs),
        "jobs": jobs,
    }


# --------------------------------------------------------------------------- #
# Shared geometry and color parsing
# --------------------------------------------------------------------------- #
def parse_color(spec, default):
    if not spec:
        return default
    spec = spec.strip()
    if spec.startswith("#"):
        h = spec.lstrip("#")
        return tuple(int(h[i : i + 2], 16) / 255 for i in (0, 2, 4))
    return tuple(float(x) for x in spec.split(","))[:3]


def parse_rect(spec):
    return tuple(float(x) for x in spec.split(","))[:4]


def parse_point(spec):
    x, y = (float(v) for v in spec.split(","))
    return (x, y)


def page_indices(a, n):
    page_spec = a.pages if hasattr(a, "pages") else None
    return parse_pages(page_spec, n) if page_spec else range(n)


def _selected_page(doc, page_number):
    if page_number < 1 or page_number > doc.page_count:
        raise PdfGoatError(f"--page must be between 1 and {doc.page_count}")
    return doc[page_number - 1]


def out_dir(a, src, suffix):
    requested_dir = a.outdir if hasattr(a, "outdir") else None
    d = Path(requested_dir or f"{Path(src).stem}_{suffix}").expanduser().resolve()
    d.mkdir(parents=True, exist_ok=True)
    return d


# --------------------------------------------------------------------------- #
# Annotate
# --------------------------------------------------------------------------- #
def cmd_annot_markup(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, a.kind))
    color = parse_color(a.color, (1, 0.9, 0) if a.kind == "highlight" else (1, 0, 0))
    doc = pymupdf.open(src)
    hits = 0
    for i in page_indices(a, doc.page_count):
        page = doc[i]
        for r in page.search_for(a.find):
            adder = {
                "highlight": page.add_highlight_annot,
                "underline": page.add_underline_annot,
                "strikeout": page.add_strikeout_annot,
            }[a.kind]
            annot = adder(r)
            annot.set_colors(stroke=color)
            annot.update()
            hits += 1
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": f"annot-{a.kind}",
        "inputs": [str(src)],
        "outputs": [out],
        "find": a.find,
        "marks": hits,
    }


def cmd_annot_note(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "note"))
    doc = pymupdf.open(src)
    page = _selected_page(doc, a.page)
    annot = page.add_text_annot(parse_point(a.at), a.text)
    annot.update()
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "annot-note",
        "inputs": [str(src)],
        "outputs": [out],
        "page": a.page,
        "text": a.text,
    }


def cmd_annot_textbox(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "textbox"))
    doc = pymupdf.open(src)
    page = _selected_page(doc, a.page)
    annot = page.add_freetext_annot(
        pymupdf.Rect(parse_rect(a.rect)),
        a.text,
        fontsize=a.size,
        text_color=parse_color(a.color, (0, 0, 0)),
        fill_color=parse_color(a.fill, (1, 1, 0.7)) if a.fill else None,
    )
    annot.update()
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "annot-textbox",
        "inputs": [str(src)],
        "outputs": [out],
        "page": a.page,
        "text": a.text,
    }


def cmd_annot_shape(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, a.kind))
    doc = pymupdf.open(src)
    page = _selected_page(doc, a.page)
    rect = pymupdf.Rect(parse_rect(a.rect))
    annot = (
        page.add_rect_annot(rect) if a.kind == "rect" else page.add_circle_annot(rect)
    )
    annot.set_colors(
        stroke=parse_color(a.color, (1, 0, 0)),
        fill=parse_color(a.fill, None) if a.fill else None,
    )
    annot.set_border(width=a.width)
    annot.update()
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": f"annot-{a.kind}",
        "inputs": [str(src)],
        "outputs": [out],
        "page": a.page,
    }


def cmd_annot_line(a):
    import pymupdf

    src = resolve(a.file)
    kind = "arrow" if a.arrow else "line"
    out = ensure_parent(a.output or default_out(src, kind))
    doc = pymupdf.open(src)
    page = _selected_page(doc, a.page)
    p1 = pymupdf.Point(parse_point(a.start))
    p2 = pymupdf.Point(parse_point(a.end))
    annot = page.add_line_annot(p1, p2)
    annot.set_colors(stroke=parse_color(a.color, (1, 0, 0)))
    annot.set_border(width=a.width)
    if a.arrow:
        annot.set_line_ends(pymupdf.PDF_ANNOT_LE_NONE, pymupdf.PDF_ANNOT_LE_OPEN_ARROW)
    annot.update()
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": f"annot-{kind}",
        "inputs": [str(src)],
        "outputs": [out],
        "page": a.page,
    }


def cmd_annot_ink(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "ink"))
    stroke = [parse_point(p) for p in a.points.split(";")]
    doc = pymupdf.open(src)
    page = _selected_page(doc, a.page)
    annot = page.add_ink_annot([stroke])
    annot.set_colors(stroke=parse_color(a.color, (0, 0, 1)))
    annot.set_border(width=a.width)
    annot.update()
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "annot-ink",
        "inputs": [str(src)],
        "outputs": [out],
        "page": a.page,
        "points": len(stroke),
    }


def cmd_annot_stamp(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "stamp"))
    doc = pymupdf.open(src)
    page = _selected_page(doc, a.page)
    annot = page.add_stamp_annot(pymupdf.Rect(parse_rect(a.rect)), stamp=a.stamp)
    annot.update()
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "annot-stamp",
        "inputs": [str(src)],
        "outputs": [out],
        "page": a.page,
        "stamp": a.stamp,
    }


def cmd_annot_callout(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "callout"))
    doc = pymupdf.open(src)
    page = _selected_page(doc, a.page)
    rect = pymupdf.Rect(parse_rect(a.rect))
    target = pymupdf.Point(parse_point(a.target))
    box = page.add_freetext_annot(
        rect, a.text, fontsize=a.size, text_color=(0, 0, 0), fill_color=(1, 1, 0.7)
    )
    box.update()
    line = page.add_line_annot(pymupdf.Point(rect.x0, rect.y1), target)
    line.set_colors(stroke=(1, 0, 0))
    line.set_line_ends(pymupdf.PDF_ANNOT_LE_NONE, pymupdf.PDF_ANNOT_LE_OPEN_ARROW)
    line.update()
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "annot-callout",
        "inputs": [str(src)],
        "outputs": [out],
        "page": a.page,
    }


def cmd_annot_area(a):
    import pymupdf

    if not 0 <= a.opacity <= 1:
        raise PdfGoatError("--opacity must be between 0 and 1")
    src = resolve(a.file)
    with pymupdf.open(src) as doc:
        if a.page < 1 or a.page > doc.page_count:
            raise PdfGoatError(f"--page must be between 1 and {doc.page_count}")
        out = ensure_parent(a.output or default_out(src, "area-highlight"))
        page = _selected_page(doc, a.page)
        annotation = page.add_rect_annot(pymupdf.Rect(parse_rect(a.rect)))
        annotation.set_colors(stroke=None, fill=parse_color(a.color, (1, 1, 0)))
        annotation.set_border(width=0)
        annotation.set_opacity(a.opacity)
        annotation.update()
        _save_pdf(doc, out)
    return {
        "verb": "annot-area-highlight",
        "inputs": [str(src)],
        "outputs": [out],
        "page": a.page,
        "opacity": a.opacity,
    }


def cmd_annot_polygon(a):
    import pymupdf

    points = [parse_point(point) for point in a.points.split(";")]
    if len(points) < 3:
        raise PdfGoatError("--points requires at least three x,y pairs")
    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "polygon"))
    doc = pymupdf.open(src)
    page = _selected_page(doc, a.page)
    annot = page.add_polygon_annot(points)
    annot.set_colors(
        stroke=parse_color(a.color, (1, 0, 0)),
        fill=parse_color(a.fill, None) if a.fill else None,
    )
    annot.set_border(width=a.width)
    annot.update()
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "annot-polygon",
        "inputs": [str(src)],
        "outputs": [out],
        "page": a.page,
        "points": len(points),
    }


def cmd_annotations(a):
    import pymupdf

    src = resolve(a.file)
    doc = pymupdf.open(src)
    items = []
    for i in range(doc.page_count):
        for annot in doc[i].annots():
            items.append(
                {
                    "page": i + 1,
                    "type": annot.type[1],
                    "rect": [round(v, 1) for v in annot.rect],
                    "content": annot.info.get("content") or None,
                    "author": annot.info.get("title") or None,
                }
            )
    doc.close()
    return {
        "verb": "annot-list",
        "inputs": [str(src)],
        "outputs": [],
        "count": len(items),
        "annotations": items,
    }


def cmd_annot_flatten(a):
    import shutil
    import subprocess

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "flat"))
    qpdf = shutil.which("qpdf")
    if not qpdf:
        raise PdfGoatError("annotate flatten requires qpdf on PATH")
    subprocess.run(
        [qpdf, "--flatten-annotations=all", str(src), out],
        check=True,
        capture_output=True,
    )
    return {"verb": "annot-flatten", "inputs": [str(src)], "outputs": [out]}


def cmd_annot_delete(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "noannots"))
    doc = pymupdf.open(src)
    removed = 0
    for i in page_indices(a, doc.page_count):
        page = doc[i]
        annot = page.first_annot
        while annot:
            if not a.type or annot.type[1].lower() == a.type.lower():
                annot = page.delete_annot(annot)
                removed += 1
            else:
                annot = annot.next
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "annot-delete",
        "inputs": [str(src)],
        "outputs": [out],
        "removed": removed,
    }


# --------------------------------------------------------------------------- #
# Forms
# --------------------------------------------------------------------------- #
def _add_widget(page, ftype, name, rect, value=None):
    import pymupdf

    w = pymupdf.Widget()
    w.field_name = name
    w.field_type = ftype
    w.rect = pymupdf.Rect(rect)
    w.border_color = (0.4, 0.4, 0.4)
    w.fill_color = (0.96, 0.96, 0.96)
    w.border_width = 1
    if value is not None:
        w.field_value = value
    return page.add_widget(w)


def cmd_form_create_text(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "field"))
    doc = pymupdf.open(src)
    _add_widget(
        _selected_page(doc, a.page),
        pymupdf.PDF_WIDGET_TYPE_TEXT,
        a.name,
        parse_rect(a.rect),
    )
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "form-create-text",
        "inputs": [str(src)],
        "outputs": [out],
        "field": a.name,
    }


def cmd_form_create_checkbox(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "checkbox"))
    doc = pymupdf.open(src)
    _add_widget(
        _selected_page(doc, a.page),
        pymupdf.PDF_WIDGET_TYPE_CHECKBOX,
        a.name,
        parse_rect(a.rect),
    )
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "form-create-checkbox",
        "inputs": [str(src)],
        "outputs": [out],
        "field": a.name,
    }


def _form_values(src):
    from pypdf import PdfReader

    fields = PdfReader(src).get_fields() or {}
    return {
        n: ("" if f.get("/V") is None else str(f.get("/V"))) for n, f in fields.items()
    }


def cmd_form_export(a):
    src = resolve(a.file)
    data = _form_values(src)
    ext = {"json": "json", "xfdf": "xfdf", "fdf": "fdf"}[a.format]
    out = ensure_parent(a.output or default_out(src, "formdata", ext))
    if a.format == "json":
        Path(out).write_text(json.dumps(data, indent=2))
    elif a.format == "xfdf":
        from xml.sax.saxutils import escape

        rows = "".join(
            f'<field name="{escape(n)}"><value>{escape(v)}</value></field>'
            for n, v in data.items()
        )
        Path(out).write_text(
            '<?xml version="1.0" encoding="UTF-8"?>\n'
            '<xfdf xmlns="http://ns.adobe.com/xfdf/"><fields>'
            f"{rows}</fields></xfdf>"
        )
    else:
        body = "".join(f"<< /T ({n}) /V ({v}) >>\n" for n, v in data.items())
        Path(out).write_text(
            "%FDF-1.2\n1 0 obj\n<< /FDF << /Fields [\n"
            + body
            + "] >> >>\nendobj\ntrailer\n"
            "<< /Root 1 0 R >>\n%%EOF\n"
        )
    return {
        "verb": "form-export",
        "inputs": [str(src)],
        "outputs": [out],
        "format": a.format,
        "field_count": len(data),
    }


def cmd_form_import(a):
    src = resolve(a.file)
    data_path = Path(a.data).expanduser()
    if data_path.suffix.lower() == ".xfdf":
        import xml.etree.ElementTree as ET

        root = ET.fromstring(data_path.read_text())
        ns = "{http://ns.adobe.com/xfdf/}"
        data = {
            f.get("name"): (f.findtext(f"{ns}value") or "")
            for f in root.iter(f"{ns}field")
        }
    else:
        data = json.loads(data_path.read_text())
    out = ensure_parent(a.output or default_out(src, "filled"))
    flattened = fill_pdf_form(src, data, out, a.flatten)
    return {
        "verb": "form-import",
        "inputs": [str(src)],
        "outputs": [out],
        "fields_set": list(data.keys()),
        "flattened": flattened,
    }


# --------------------------------------------------------------------------- #
# Security
# --------------------------------------------------------------------------- #
def cmd_sec_encrypt(a):
    import pikepdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "encrypted"))
    enc = pikepdf.Encryption(owner=a.owner or a.password, user=a.password, R=6)
    with pikepdf.open(src) as pdf:
        pdf.save(out, encryption=enc)
    return {
        "verb": "sec-encrypt",
        "inputs": [str(src)],
        "outputs": [out],
        "algorithm": "AES-256",
    }


def cmd_sec_decrypt(a):
    import pikepdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "decrypted"))
    with pikepdf.open(src, password=a.password) as pdf:
        pdf.save(out)
    return {"verb": "sec-decrypt", "inputs": [str(src)], "outputs": [out]}


def cmd_sec_permissions(a):
    import pikepdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "restricted"))
    allow = pikepdf.Permissions(
        extract=not a.no_copy,
        modify_annotation=not a.no_modify,
        modify_assembly=not a.no_modify,
        modify_form=not a.no_modify,
        modify_other=not a.no_modify,
        print_lowres=not a.no_print,
        print_highres=not a.no_print,
    )
    enc = pikepdf.Encryption(owner=a.owner, user=a.user or "", allow=allow, R=6)
    with pikepdf.open(src) as pdf:
        pdf.save(out, encryption=enc)
    return {
        "verb": "sec-permissions",
        "inputs": [str(src)],
        "outputs": [out],
        "no_print": a.no_print,
        "no_copy": a.no_copy,
        "no_modify": a.no_modify,
    }


def _self_signed(cn):
    import datetime

    from cryptography import x509
    from cryptography.hazmat.primitives import hashes, serialization
    from cryptography.hazmat.primitives.asymmetric import rsa
    from cryptography.x509.oid import NameOID

    key = rsa.generate_private_key(public_exponent=65537, key_size=2048)
    name = x509.Name([x509.NameAttribute(NameOID.COMMON_NAME, cn)])
    now = datetime.datetime.now(datetime.timezone.utc)
    cert = (
        x509.CertificateBuilder()
        .subject_name(name)
        .issuer_name(name)
        .public_key(key.public_key())
        .serial_number(x509.random_serial_number())
        .not_valid_before(now - datetime.timedelta(days=1))
        .not_valid_after(now + datetime.timedelta(days=3650))
        .add_extension(x509.BasicConstraints(ca=False, path_length=None), critical=True)
        .add_extension(
            x509.KeyUsage(
                digital_signature=True,
                content_commitment=True,
                key_encipherment=False,
                data_encipherment=False,
                key_agreement=False,
                key_cert_sign=False,
                crl_sign=False,
                encipher_only=False,
                decipher_only=False,
            ),
            critical=True,
        )
        .sign(key, hashes.SHA256())
    )
    key_pem = key.private_bytes(
        serialization.Encoding.PEM,
        serialization.PrivateFormat.PKCS8,
        serialization.NoEncryption(),
    )
    return key_pem, cert.public_bytes(serialization.Encoding.PEM)


def cmd_sec_sign(a):
    import tempfile

    from pyhanko.pdf_utils.incremental_writer import IncrementalPdfFileWriter
    from pyhanko.sign import signers

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "signed"))
    cn = a.name or "pdf-goat demo"
    key_pem, cert_pem = _self_signed(cn)
    with tempfile.TemporaryDirectory() as td:
        kp, cp = Path(td) / "k.pem", Path(td) / "c.pem"
        kp.write_bytes(key_pem)
        cp.write_bytes(cert_pem)
        signer = signers.SimpleSigner.load(str(kp), str(cp))
        with open(src, "rb") as inf, open(out, "wb") as outf:
            w = IncrementalPdfFileWriter(inf)
            meta = signers.PdfSignatureMetadata(
                field_name=a.field or "Signature1",
                reason=a.reason or "Approval",
                name=cn,
            )
            signers.sign_pdf(w, meta, signer=signer, output=outf)
    return {
        "verb": "sec-sign",
        "inputs": [str(src)],
        "outputs": [out],
        "signer": cn,
        "self_signed": True,
    }


def cmd_sec_verify(a):
    from pyhanko.pdf_utils.reader import PdfFileReader
    from pyhanko.sign.validation import validate_pdf_signature
    from pyhanko_certvalidator import ValidationContext

    src = resolve(a.file)
    sigs = []
    with open(src, "rb") as f:
        reader = PdfFileReader(f)
        for sig in reader.embedded_signatures:
            entry = {
                "field": sig.field_name,
                "signer": sig.signer_cert.subject.human_friendly,
            }
            try:
                vc = ValidationContext(
                    trust_roots=[sig.signer_cert], allow_fetching=False
                )
                status = validate_pdf_signature(sig, vc)
                entry.update(
                    intact=bool(status.intact),
                    valid=bool(status.valid),
                    trusted=bool(status.trusted),
                    coverage=str(
                        status.coverage if hasattr(status, "coverage") else ""
                    ),
                    modified=(
                        status.modification_level.name
                        if hasattr(status, "modification_level")
                        and status.modification_level
                        else None
                    ),
                )
            except Exception as e:  # noqa: BLE001
                entry["validation_error"] = str(e)
            sigs.append(entry)
    return {
        "verb": "sec-verify",
        "inputs": [str(src)],
        "outputs": [],
        "signature_count": len(sigs),
        "signatures": sigs,
    }


def cmd_sec_sanitize(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "sanitized"))
    doc = pymupdf.open(src)
    doc.scrub(
        attached_files=True,
        embedded_files=True,
        javascript=True,
        xml_metadata=True,
        thumbnails=True,
        clean_pages=True,
        hidden_text=False,
        metadata=False,
        remove_links=False,
        reset_fields=False,
        reset_responses=False,
        redactions=False,
    )
    _save_pdf(doc, out, clean=True)
    doc.close()
    return {
        "verb": "sec-sanitize",
        "inputs": [str(src)],
        "outputs": [out],
        "removed": [
            "javascript",
            "embedded_files",
            "attached_files",
            "xml_metadata",
            "thumbnails",
        ],
    }


# --------------------------------------------------------------------------- #
# Metadata
# --------------------------------------------------------------------------- #
def cmd_meta_get(a):
    import pymupdf

    src = resolve(a.file)
    doc = pymupdf.open(src)
    md = dict(doc.metadata)
    has_xml = bool(doc.xref_xml_metadata if hasattr(doc, "xref_xml_metadata") else 0)
    doc.close()
    return {
        "verb": "meta-get",
        "inputs": [str(src)],
        "outputs": [],
        "metadata": {k: v for k, v in md.items() if v},
        "has_xmp": has_xml,
    }


def cmd_meta_set(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "meta"))
    doc = pymupdf.open(src)
    md = dict(doc.metadata)
    for item in a.set or []:
        k, _, v = item.partition("=")
        md[k.strip()] = v
    doc.set_metadata(md)
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "meta-set",
        "inputs": [str(src)],
        "outputs": [out],
        "metadata": {k: v for k, v in md.items() if v},
    }


def cmd_meta_strip(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "nometa"))
    doc = pymupdf.open(src)
    doc.set_metadata({})
    doc.del_xml_metadata()
    _save_pdf(doc, out)
    doc.close()
    return {"verb": "meta-strip", "inputs": [str(src)], "outputs": [out]}


# --------------------------------------------------------------------------- #
# Page layout and numbering
# --------------------------------------------------------------------------- #
def cmd_pages_blank(a):
    import pymupdf

    if a.count < 1 or a.count > 100:
        raise PdfGoatError("--count must be between 1 and 100")
    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "blank-pages"))
    doc = pymupdf.open(src)
    max_at = doc.page_count + 1
    at = a.at if a.at is not None else max_at
    if at < 1 or at > max_at:
        doc.close()
        raise PdfGoatError(f"--at must be between 1 and {max_at}")
    template_index = min(max(at - 2, 0), doc.page_count - 1)
    template = doc[template_index].rect
    insertion_index = at - 1
    for offset in range(a.count):
        doc.new_page(
            pno=insertion_index + offset,
            width=template.width,
            height=template.height,
        )
    _save_pdf(doc, out)
    total_pages = doc.page_count
    doc.close()
    return {
        "verb": "pages-blank",
        "inputs": [str(src)],
        "outputs": [out],
        "inserted_at": at,
        "inserted_pages": a.count,
        "total_pages": total_pages,
    }


def cmd_pages_duplicate(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "duplicated"))
    doc = pymupdf.open(src)
    donor = pymupdf.open(src)
    selected = sorted(set(parse_pages(a.pages, doc.page_count)))
    for offset, page_index in enumerate(selected):
        doc.insert_pdf(
            donor,
            from_page=page_index,
            to_page=page_index,
            start_at=page_index + offset + 1,
        )
    _save_pdf(doc, out)
    total_pages = doc.page_count
    donor.close()
    doc.close()
    return {
        "verb": "pages-duplicate",
        "inputs": [str(src)],
        "outputs": [out],
        "duplicated_pages": [index + 1 for index in selected],
        "total_pages": total_pages,
    }


def cmd_pages_crop(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "cropped"))
    box = pymupdf.Rect(parse_rect(a.box))
    doc = pymupdf.open(src)
    for i in page_indices(a, doc.page_count):
        doc[i].set_cropbox(box)
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "pages-crop",
        "inputs": [str(src)],
        "outputs": [out],
        "box": list(box),
    }


def cmd_pages_scale(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "scaled"))
    doc = pymupdf.open(src)
    new = pymupdf.open()
    for i in range(doc.page_count):
        r = doc[i].rect
        page = new.new_page(width=r.width * a.factor, height=r.height * a.factor)
        page.show_pdf_page(page.rect, doc, i)
    _save_pdf(new, out)
    return {
        "verb": "pages-scale",
        "inputs": [str(src)],
        "outputs": [out],
        "factor": a.factor,
    }


def cmd_pages_nup(a):
    import pymupdf

    src = resolve(a.file)
    cols, rows = (2, 1) if a.n == 2 else (2, 2)
    out = ensure_parent(a.output or default_out(src, f"{a.n}up"))
    doc = pymupdf.open(src)
    w, h = doc[0].rect.width, doc[0].rect.height
    new = pymupdf.open()
    per = cols * rows
    for start in range(0, doc.page_count, per):
        sheet = new.new_page(width=cols * w, height=rows * h)
        for k in range(per):
            pno = start + k
            if pno >= doc.page_count:
                break
            c, r = k % cols, k // cols
            sheet.show_pdf_page(
                pymupdf.Rect(c * w, r * h, (c + 1) * w, (r + 1) * h), doc, pno
            )
    _save_pdf(new, out)
    return {
        "verb": "pages-nup",
        "inputs": [str(src)],
        "outputs": [out],
        "n": a.n,
        "sheets": new.page_count,
    }


def cmd_pages_booklet(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "booklet"))
    doc = pymupdf.open(src)
    n = doc.page_count
    padded = n + (-n % 4)
    order = []
    lo, hi = 0, padded - 1
    while lo < hi:
        order += [hi, lo, lo + 1, hi - 1]
        lo += 2
        hi -= 2
    w, h = doc[0].rect.width, doc[0].rect.height
    new = pymupdf.open()
    for k in range(0, len(order), 2):
        sheet = new.new_page(width=2 * w, height=h)
        for slot in (0, 1):
            pno = order[k + slot]
            if pno < n:
                sheet.show_pdf_page(
                    pymupdf.Rect(slot * w, 0, (slot + 1) * w, h), doc, pno
                )
    _save_pdf(new, out)
    return {
        "verb": "pages-booklet",
        "inputs": [str(src)],
        "outputs": [out],
        "sheets": new.page_count,
        "page_order": [p + 1 for p in order],
    }


def _stamp_text(doc, text_fn, where, fontsize, color):
    import pymupdf

    for i in range(doc.page_count):
        page = doc[i]
        text = text_fn(i, doc.page_count)
        tl = pymupdf.get_text_length(text, fontsize=fontsize)
        r = page.rect
        x = {"left": 40, "center": (r.width - tl) / 2, "right": r.width - tl - 40}[
            where[1]
        ]
        y = 50 if where[0] == "top" else r.height - 36
        page.insert_text((x, y), text, fontsize=fontsize, color=color)


def cmd_pages_header(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, a.where))
    doc = pymupdf.open(src)

    def fn(i, n):
        return a.text.replace("{page}", str(i + 1)).replace("{pages}", str(n))

    vpos = "top" if a.where == "header" else "bottom"
    _stamp_text(doc, fn, (vpos, a.align), a.size, parse_color(a.color, (0.2, 0.2, 0.2)))
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": f"pages-{a.where}",
        "inputs": [str(src)],
        "outputs": [out],
        "text": a.text,
    }


def cmd_pages_numbers(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "numbered"))
    doc = pymupdf.open(src)

    def fn(i, n):
        return a.format.replace("{page}", str(i + a.start)).replace("{pages}", str(n))

    _stamp_text(doc, fn, ("bottom", a.align), a.size, (0.2, 0.2, 0.2))
    _save_pdf(doc, out)
    doc.close()
    return {"verb": "pages-numbers", "inputs": [str(src)], "outputs": [out]}


def cmd_pages_bates(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "bates"))
    doc = pymupdf.open(src)

    def fn(i, n):
        return f"{a.prefix}{str(a.start + i).zfill(a.digits)}"

    _stamp_text(doc, fn, ("bottom", "right"), a.size, (0.1, 0.1, 0.1))
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "pages-bates",
        "inputs": [str(src)],
        "outputs": [out],
        "first": f"{a.prefix}{str(a.start).zfill(a.digits)}",
    }


def cmd_pages_boxes(a):
    import pikepdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "boxes"))
    key = {
        "media": "/MediaBox",
        "crop": "/CropBox",
        "trim": "/TrimBox",
        "bleed": "/BleedBox",
    }[a.box]
    rect = list(parse_rect(a.rect))
    with pikepdf.open(src) as pdf:
        for i in page_indices(a, len(pdf.pages)):
            pdf.pages[i][key] = rect
        pdf.save(out)
    return {
        "verb": "pages-boxes",
        "inputs": [str(src)],
        "outputs": [out],
        "box": a.box,
        "rect": rect,
    }


def cmd_pages_insert(a):
    import pikepdf

    src = resolve(a.file)
    other = resolve(a.source)
    out = ensure_parent(a.output or default_out(src, "inserted"))
    with pikepdf.open(src) as pdf, pikepdf.open(other) as ins:
        at = (a.at - 1) if a.at else len(pdf.pages)
        for off, p in enumerate(ins.pages):
            pdf.pages.insert(at + off, p)
        pdf.save(out)
    return {
        "verb": "pages-insert",
        "inputs": [str(src), str(other)],
        "outputs": [out],
        "at": a.at,
    }


def cmd_pages_replace(a):
    import pikepdf

    src = resolve(a.file)
    other = resolve(a.source)
    out = ensure_parent(a.output or default_out(src, "replaced"))
    with pikepdf.open(src) as pdf, pikepdf.open(other) as rep:
        targets = sorted(set(parse_pages(a.pages, len(pdf.pages))))
        at = targets[0]
        for i in reversed(targets):
            del pdf.pages[i]
        for off, p in enumerate(rep.pages):
            pdf.pages.insert(at + off, p)
        pdf.save(out)
    return {
        "verb": "pages-replace",
        "inputs": [str(src), str(other)],
        "outputs": [out],
        "replaced": [t + 1 for t in targets],
    }


# --------------------------------------------------------------------------- #
# Extract assets
# --------------------------------------------------------------------------- #
def _extract_image_with_mupdf(doc, xref, outdir):
    """Write an image pikepdf could not extract.

    pikepdf refuses a CMYK JPEG with an Adobe colour transform, and MuPDF
    would re-encode that one, so a lone DCT stream with no Decode array is
    copied as stored. Everything else goes through MuPDF.
    """
    lone_dct = doc.xref_get_key(xref, "Filter") in (
        ("name", "/DCTDecode"),
        ("array", "[/DCTDecode]"),
    )
    if lone_dct and doc.xref_get_key(xref, "Decode") == ("null", "null"):
        written = outdir / f"img_{xref}.jpg"
        written.write_bytes(doc.xref_stream_raw(xref))
    else:
        extracted = doc.extract_image(xref)
        written = outdir / f"img_{xref}.{extracted['ext']}"
        written.write_bytes(extracted["image"])
    return str(written)


def cmd_get_images(a):
    import pikepdf
    import pymupdf

    src = resolve(a.file)
    outdir = out_dir(a, src, "images")
    outputs, seen = [], set()
    # MuPDF finds the images, including those nested in form XObjects; pikepdf
    # writes the compressed stream out as it stands instead of re-encoding it.
    # pikepdf declines with a dozen exception types (its own, Pillow's,
    # NotImplementedError, ValueError), and MuPDF reads all of them.
    with pymupdf.open(src) as doc, pikepdf.open(src) as pdf:
        for i in range(doc.page_count):
            for img in doc.get_page_images(i):
                xref = img[0]
                if xref in seen:
                    continue
                seen.add(xref)
                obj = pdf.get_object(xref, 0)
                if obj is None:  # the object sits at a later generation
                    written = _extract_image_with_mupdf(doc, xref, outdir)
                else:
                    try:
                        written = pikepdf.PdfImage(obj).extract_to(
                            fileprefix=str(outdir / f"img_{xref}")
                        )
                    except Exception:  # noqa: BLE001
                        written = _extract_image_with_mupdf(doc, xref, outdir)
                outputs.append(written)
    return {
        "verb": "get-images",
        "inputs": [str(src)],
        "outputs": outputs,
        "count": len(outputs),
    }


def cmd_get_fonts(a):
    import pymupdf

    src = resolve(a.file)
    doc = pymupdf.open(src)
    fonts = {}
    for i in range(doc.page_count):
        for f in doc.get_page_fonts(i):
            fonts[f[0]] = {"name": f[3], "type": f[2], "ext": f[1], "encoding": f[5]}
    doc.close()
    return {
        "verb": "get-fonts",
        "inputs": [str(src)],
        "outputs": [],
        "count": len(fonts),
        "fonts": list(fonts.values()),
    }


def _iter_text_blocks(doc, indices):
    for page_index in indices:
        page = doc[page_index]
        for block in page.get_text("blocks", sort=True):
            x0, y0, x1, y1, text, block_number, block_type, *_ = block
            if block_type != 0 or not text.strip():
                continue
            yield {
                "page": page_index + 1,
                "block": block_number,
                "rect": [round(value, 2) for value in (x0, y0, x1, y1)],
                "text": text.rstrip(),
            }


def cmd_get_text_blocks(a):
    import pymupdf

    if a.max_blocks < 1 or a.max_blocks > 1000:
        raise PdfGoatError("--max-blocks must be between 1 and 1000")
    if a.start_block < 0:
        raise PdfGoatError("--start-block must be zero or greater")
    src = resolve(a.file)
    with pymupdf.open(src) as doc:
        total_pages = doc.page_count
        page_truncated = not a.pages and total_pages > _DEFAULT_PAGE_WINDOW
        indices = (
            parse_pages(a.pages, total_pages)
            if a.pages
            else range(min(_DEFAULT_PAGE_WINDOW, total_pages))
        )
        captured = list(
            islice(
                _iter_text_blocks(doc, indices),
                a.start_block,
                a.start_block + a.max_blocks + 1,
            )
        )
    block_truncated = len(captured) > a.max_blocks
    blocks = captured[: a.max_blocks]
    return {
        "verb": "get-text-blocks",
        "inputs": [str(src)],
        "outputs": [],
        "total_pages": total_pages,
        "count": len(blocks),
        "blocks": blocks,
        "start_block": a.start_block,
        "page_truncated": page_truncated,
        "block_truncated": block_truncated,
        "next_block": a.start_block + a.max_blocks if block_truncated else None,
        "next_page": _DEFAULT_PAGE_WINDOW + 1
        if page_truncated and not block_truncated
        else None,
        "truncated": page_truncated or block_truncated,
    }


def cmd_get_attachments(a):
    import pymupdf

    src = resolve(a.file)
    outdir = Path(a.outdir or f"{src.stem}_attachments").expanduser().resolve()
    with pymupdf.open(src) as doc:
        entries = [(name, "catalog", name) for name in doc.embfile_names()]
        entries.extend(
            (name, "annotation", (page, annotation))
            for page, annotation, name in _file_attachment_annotations(doc)
        )
        planned = []
        destination_keys = set()
        for name, kind, source in entries:
            safe_name = Path(name).name
            if safe_name in {"", ".", ".."}:
                raise PdfGoatError("an attachment has no safe file name")
            destination = outdir / safe_name
            destination_key = unicodedata.normalize("NFC", safe_name).casefold()
            if (
                destination_key in destination_keys
                or destination.exists()
                or destination.is_symlink()
            ):
                raise PdfGoatError(f"attachment output already exists: {destination}")
            destination_keys.add(destination_key)
            planned.append((destination, kind, source))
        outdir.mkdir(parents=True, exist_ok=True)

        outputs = []
        for destination, kind, source in planned:
            payload = (
                doc.embfile_get(source) if kind == "catalog" else source[1].get_file()
            )
            try:
                with destination.open("xb") as output:
                    output.write(payload)
            except FileExistsError as error:
                raise PdfGoatError(
                    f"attachment output already exists: {destination}"
                ) from error
            outputs.append(str(destination))
    return {
        "verb": "get-attachments",
        "inputs": [str(src)],
        "outputs": outputs,
        "count": len(outputs),
    }


def cmd_get_bookmarks(a):
    import pymupdf

    src = resolve(a.file)
    doc = pymupdf.open(src)
    toc = [
        {"level": lvl, "title": title, "page": page}
        for lvl, title, page in doc.get_toc()
    ]
    doc.close()
    return {
        "verb": "get-bookmarks",
        "inputs": [str(src)],
        "outputs": [],
        "count": len(toc),
        "bookmarks": toc,
    }


def cmd_get_links(a):
    import pymupdf

    src = resolve(a.file)
    doc = pymupdf.open(src)
    links = []
    for i in range(doc.page_count):
        for lk in doc[i].get_links():
            links.append(
                {
                    "page": i + 1,
                    "rect": [round(v, 1) for v in lk["from"]],
                    "uri": lk.get("uri"),
                    "target_page": lk.get("page"),
                }
            )
    doc.close()
    return {
        "verb": "get-links",
        "inputs": [str(src)],
        "outputs": [],
        "count": len(links),
        "links": links,
    }


# --------------------------------------------------------------------------- #
# Outline and link editing
# --------------------------------------------------------------------------- #
def cmd_bookmarks_set(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "outlined"))
    toc = json.loads(Path(a.data).expanduser().read_text())
    doc = pymupdf.open(src)
    doc.set_toc([[e["level"], e["title"], e["page"]] for e in toc])
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "bookmarks-set",
        "inputs": [str(src)],
        "outputs": [out],
        "count": len(toc),
    }


def cmd_bookmarks_clear(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "unbookmarked"))
    doc = pymupdf.open(src)
    removed = len(doc.get_toc())
    doc.set_toc([])
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "bookmarks-clear",
        "inputs": [str(src)],
        "outputs": [out],
        "removed": removed,
    }


def cmd_links_add(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "linked"))
    doc = pymupdf.open(src)
    page = _selected_page(doc, a.page)
    rect = pymupdf.Rect(parse_rect(a.rect))
    if a.uri:
        page.insert_link({"kind": pymupdf.LINK_URI, "from": rect, "uri": a.uri})
    else:
        page.insert_link({"kind": pymupdf.LINK_GOTO, "from": rect, "page": a.goto - 1})
    _save_pdf(doc, out)
    doc.close()
    return {"verb": "links-add", "inputs": [str(src)], "outputs": [out], "page": a.page}


def cmd_links_remove(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "links-removed"))
    doc = pymupdf.open(src)
    removed = 0
    for page_index in page_indices(a, doc.page_count):
        page = doc[page_index]
        for link in page.get_links():
            if a.external_only and link["kind"] not in _external_link_kinds():
                continue
            page.delete_link(link)
            removed += 1
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "links-remove",
        "inputs": [str(src)],
        "outputs": [out],
        "removed": removed,
        "external_only": a.external_only,
    }


# --------------------------------------------------------------------------- #
# Convert
# --------------------------------------------------------------------------- #
def cmd_convert_ocr(a):
    import ocrmypdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "ocr"))
    ocrmypdf.ocr(
        str(src),
        out,
        force_ocr=a.force,
        skip_text=not a.force,
        progress_bar=False,
        deskew=False,
    )
    return {"verb": "convert-ocr", "inputs": [str(src)], "outputs": [out]}


def cmd_convert_html(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "from-pdf", "html"))
    doc = pymupdf.open(src)
    parts = [doc[i].get_text("html") for i in range(doc.page_count)]
    doc.close()
    Path(out).write_text(
        "<!DOCTYPE html><html><head><meta charset='utf-8'></head><body>"
        + "\n<hr/>\n".join(parts)
        + "</body></html>"
    )
    return {"verb": "convert-html", "inputs": [str(src)], "outputs": [out]}


def cmd_convert_tables(a):
    import csv

    import pdfplumber

    src = resolve(a.file)
    outdir = out_dir(a, src, "tables")
    outputs = []
    with pdfplumber.open(src) as pdf:
        for i, page in enumerate(pdf.pages):
            for t, table in enumerate(page.extract_tables()):
                p = outdir / f"p{i + 1}_t{t + 1}.csv"
                with open(p, "w", newline="") as fh:
                    csv.writer(fh).writerows(table)
                outputs.append(str(p))
    return {
        "verb": "convert-tables",
        "inputs": [str(src)],
        "outputs": outputs,
        "tables": len(outputs),
    }


def cmd_convert_pdfa(a):
    import shutil
    import subprocess

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "pdfa"))
    gs = shutil.which("gs")
    if not gs:
        raise PdfGoatError("convert pdfa requires Ghostscript (gs) on PATH")
    proc = subprocess.run(
        [
            gs,
            "-dPDFA=2",
            "-dBATCH",
            "-dNOPAUSE",
            "-dQUIET",
            "-sDEVICE=pdfwrite",
            "-sColorConversionStrategy=RGB",
            "-dPDFACompatibilityPolicy=1",
            f"-sOutputFile={out}",
            str(src),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    if proc.returncode != 0 or not Path(out).exists():
        raise PdfGoatError(f"ghostscript pdf/a failed: {proc.stderr.strip()[:200]}")
    return {
        "verb": "convert-pdfa",
        "inputs": [str(src)],
        "outputs": [out],
        "standard": "PDF/A-2b",
        "note": "PDF/A conformance was not validated.",
    }


# --------------------------------------------------------------------------- #
# Office conversion
# --------------------------------------------------------------------------- #
def cmd_convert_docx(a):
    from pdf2docx import Converter

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "from-pdf", "docx"))
    cv = Converter(str(src))
    try:
        cv.convert(out, start=0, end=None)
    finally:
        cv.close()
    return {"verb": "convert-docx", "inputs": [str(src)], "outputs": [out]}


def cmd_convert_xlsx(a):
    import pdfplumber
    from openpyxl import Workbook

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "from-pdf", "xlsx"))
    wb = Workbook()
    wb.remove(wb.active)
    sheets = 0
    with pdfplumber.open(src) as pdf:
        for i, page in enumerate(pdf.pages):
            tables = page.extract_tables()
            if not tables:
                continue
            ws = wb.create_sheet(title=f"page{i + 1}"[:31])
            sheets += 1
            for table in tables:
                for row in table:
                    ws.append([("" if c is None else c) for c in row])
                ws.append([])
    if sheets == 0:
        wb.create_sheet(title="empty")
    wb.save(out)
    return {
        "verb": "convert-xlsx",
        "inputs": [str(src)],
        "outputs": [out],
        "sheets": sheets,
    }


def cmd_convert_pptx(a):
    import pymupdf
    from pptx import Presentation
    from pptx.util import Emu

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "from-pdf", "pptx"))
    import tempfile

    doc = pymupdf.open(src)
    prs = Presentation()
    blank = prs.slide_layouts[6]
    with tempfile.TemporaryDirectory() as td:
        for i in range(doc.page_count):
            r = doc[i].rect
            prs.slide_width = Emu(int(r.width / 72 * 914400))
            prs.slide_height = Emu(int(r.height / 72 * 914400))
            png = Path(td) / f"p{i}.png"
            doc[i].get_pixmap(dpi=a.dpi).save(png)
            slide = prs.slides.add_slide(blank)
            slide.shapes.add_picture(str(png), 0, 0, prs.slide_width, prs.slide_height)
    doc.close()
    prs.save(out)
    return {
        "verb": "convert-pptx",
        "inputs": [str(src)],
        "outputs": [out],
        "slides": len(prs.slides._sldIdLst),
    }


def cmd_convert_from_office(a):
    import shutil
    import subprocess

    import pymupdf

    src = resolve(a.file)
    office2pdf = shutil.which("office2pdf")
    if not office2pdf:
        raise PdfGoatError(
            "convert from-office requires office2pdf on PATH "
            "(cargo install office2pdf-cli)"
        )
    out = ensure_parent(a.output or str(src.with_suffix(".pdf")))
    # office2pdf truncates its output path before writing, so it writes a
    # sibling that only replaces `out` once the run has succeeded.
    with AtomicOutput(out) as partial:
        try:
            proc = subprocess.run(
                [office2pdf, str(src), "-o", partial],
                capture_output=True,
                text=True,
                timeout=180,
                check=False,
            )
        except subprocess.TimeoutExpired:
            raise PdfGoatError("office2pdf timed out after 180s")
        if proc.returncode != 0 or not Path(partial).exists():
            raise PdfGoatError(
                f"office2pdf failed: {(proc.stderr or proc.stdout).strip()[:200]}"
            )
    warnings = [
        line.removeprefix("Warning: ")
        for line in proc.stderr.splitlines()
        if line.startswith("Warning: ")
    ]
    with pymupdf.open(out) as doc:
        pages = doc.page_count
    return {
        "verb": "convert-from-office",
        "inputs": [str(src)],
        "outputs": [out],
        "engine": "office2pdf",
        "pages": pages,
        "warnings": warnings,
    }


def cmd_convert_audio(a):
    import shutil
    import subprocess

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "audio", "aiff"))
    say = shutil.which("say")
    if not say:
        raise PdfGoatError("convert audio requires macOS say on PATH")
    import pymupdf

    doc = pymupdf.open(src)
    text = "\n".join(doc[i].get_text("text") for i in range(doc.page_count)).strip()
    doc.close()
    if not text:
        raise PdfGoatError("no extractable text to read aloud")
    import tempfile

    with tempfile.NamedTemporaryFile("w", suffix=".txt", delete=False) as tf:
        tf.write(text)
        txt_path = tf.name
    cmd = [say, "-o", out, "-f", txt_path]
    if a.voice:
        cmd += ["-v", a.voice]
    proc = subprocess.run(cmd, capture_output=True, text=True, check=False)
    Path(txt_path).unlink(missing_ok=True)
    if proc.returncode != 0 or not Path(out).exists():
        raise PdfGoatError(f"say failed: {proc.stderr.strip()[:200]}")
    duration = None
    afinfo = shutil.which("afinfo")
    if afinfo:
        info = subprocess.run(
            [afinfo, out], capture_output=True, text=True, check=False
        ).stdout
        m = re.search(r"estimated duration:\s*([\d.]+)", info)
        if m:
            duration = round(float(m.group(1)), 2)
    return {
        "verb": "convert-audio",
        "inputs": [str(src)],
        "outputs": [out],
        "chars": len(text),
        "voice": a.voice,
        "duration_sec": duration,
    }


# --------------------------------------------------------------------------- #
# Optimize
# --------------------------------------------------------------------------- #
def cmd_optimize_reduce(a):
    import shutil
    import subprocess

    import pikepdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "reduced"))
    orig = src.stat().st_size
    gs = shutil.which("gs")
    if not gs:
        raise PdfGoatError("optimize reduce requires Ghostscript (gs) on PATH")
    proc = subprocess.run(
        [
            gs,
            "-sDEVICE=pdfwrite",
            f"-dPDFSETTINGS=/{a.preset}",
            "-dCompatibilityLevel=1.5",
            "-dNOPAUSE",
            "-dBATCH",
            "-dQUIET",
            f"-sOutputFile={out}",
            str(src),
        ],
        capture_output=True,
        text=True,
        check=False,
    )
    if proc.returncode != 0 or not Path(out).exists():
        raise PdfGoatError(f"ghostscript reduce failed: {proc.stderr.strip()[:200]}")
    with pikepdf.open(out, allow_overwriting_input=True) as pdf:
        pdf.save(out, linearize=True)
    new = Path(out).stat().st_size
    if new >= orig:
        shutil.copyfile(src, out)
        new = orig
    return {
        "verb": "optimize-reduce",
        "inputs": [str(src)],
        "outputs": [out],
        "preset": a.preset,
        "original_bytes": orig,
        "reduced_bytes": new,
        "ratio": round(new / orig, 3) if orig else None,
        "saved_bytes": orig - new,
    }


# --------------------------------------------------------------------------- #
# Text replacement within existing runs
# --------------------------------------------------------------------------- #
def cmd_edit_text(a):
    import pymupdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "edited"))
    doc = pymupdf.open(src)
    count = 0
    for page in doc:
        pending = []
        for block in page.get_text("dict")["blocks"]:
            for line in block.get("lines", []):
                for span in line["spans"]:
                    if a.find in span["text"]:
                        pending.append(
                            (
                                pymupdf.Rect(span["bbox"]),
                                span["text"].replace(a.find, a.replace),
                                span["size"],
                                pymupdf.sRGB_to_pdf(span["color"]),
                                span["origin"],
                            )
                        )
                        page.add_redact_annot(
                            pymupdf.Rect(span["bbox"]), fill=(1, 1, 1)
                        )
        if pending:
            page.apply_redactions()
            for rect, new_text, size, color, origin in pending:
                page.insert_text(origin, new_text, fontsize=size, color=color)
                count += 1
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "edit-text",
        "inputs": [str(src)],
        "outputs": [out],
        "find": a.find,
        "replace": a.replace,
        "replacements": count,
        "note": "Replaces simple text runs only. It does not reflow text or match embedded fonts.",
    }


# --------------------------------------------------------------------------- #
# Accessibility
# --------------------------------------------------------------------------- #
def cmd_access_check(a):
    import pikepdf
    import pymupdf

    src = resolve(a.file)
    with pikepdf.open(src) as pdf:
        root = pdf.Root
        tagged = "/StructTreeRoot" in root
        marked = (
            bool(root.MarkInfo.Marked)
            if "/MarkInfo" in root and "/Marked" in root.MarkInfo
            else False
        )
        lang = str(root.Lang) if "/Lang" in root else None
        title = None
        if pdf.docinfo is not None and "/Title" in pdf.docinfo:
            title = str(pdf.docinfo.Title) or None
        alt = 0
        for obj in pdf.objects:
            if isinstance(obj, pikepdf.Dictionary) and "/Alt" in obj:
                alt += 1
    doc = pymupdf.open(src)
    images = sum(len(doc.get_page_images(i)) for i in range(doc.page_count))
    doc.close()
    return {
        "verb": "access-check",
        "inputs": [str(src)],
        "outputs": [],
        "tagged": tagged or marked,
        "has_title": bool(title),
        "title": title,
        "has_lang": bool(lang),
        "lang": lang,
        "images": images,
        "images_without_alt": max(0, images - alt),
        "issues": [
            k
            for k, ok in [
                ("untagged", not (tagged or marked)),
                ("no_title", not title),
                ("no_lang", not lang),
                ("images_missing_alt", images - alt > 0),
            ]
            if ok
        ],
    }


def cmd_access_set(a):
    import pikepdf

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "accessible"))
    with pikepdf.open(src) as pdf:
        if a.lang:
            pdf.Root.Lang = pikepdf.String(a.lang)
        pdf.Root.MarkInfo = pikepdf.Dictionary(Marked=True)
        if a.title:
            with pdf.open_metadata() as meta:
                meta["dc:title"] = a.title
            pdf.docinfo["/Title"] = a.title
        pdf.save(out)
    return {
        "verb": "access-set",
        "inputs": [str(src)],
        "outputs": [out],
        "title": a.title,
        "lang": a.lang,
        "note": "Sets the Marked, Lang, and Title entries. It does not create a full tag tree.",
    }


# --------------------------------------------------------------------------- #
# Flatten annotations, forms, and transparency
# --------------------------------------------------------------------------- #
def cmd_pages_flatten(a):
    import shutil
    import subprocess

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "flattened"))
    qpdf = shutil.which("qpdf")
    gs = shutil.which("gs")
    if not qpdf:
        raise PdfGoatError("pages flatten requires qpdf on PATH")
    import tempfile

    tmp = Path(tempfile.gettempdir()) / f"{src.stem}.flat1.pdf"
    subprocess.run(
        [
            qpdf,
            "--generate-appearances",
            "--flatten-annotations=all",
            str(src),
            str(tmp),
        ],
        check=True,
        capture_output=True,
    )
    flattened_transparency = False
    if gs:
        proc = subprocess.run(
            [
                gs,
                "-sDEVICE=pdfwrite",
                "-dCompatibilityLevel=1.4",
                "-dNOPAUSE",
                "-dBATCH",
                "-dQUIET",
                "-dPreserveAnnots=false",
                f"-sOutputFile={out}",
                str(tmp),
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        flattened_transparency = proc.returncode == 0 and Path(out).exists()
    if not flattened_transparency:
        shutil.move(str(tmp), out)
    else:
        tmp.unlink(missing_ok=True)
    return {
        "verb": "pages-flatten",
        "inputs": [str(src)],
        "outputs": [out],
        "flattened": ["annotations", "form_fields"]
        + (["transparency"] if flattened_transparency else []),
    }


# --------------------------------------------------------------------------- #
# Compare
# --------------------------------------------------------------------------- #
def cmd_compare_text(a):
    import difflib

    import pymupdf

    da, db = pymupdf.open(resolve(a.file)), pymupdf.open(resolve(a.other))
    ta = "\n".join(p.get_text("text") for p in da).splitlines()
    tb = "\n".join(p.get_text("text") for p in db).splitlines()
    da.close()
    db.close()
    diff = list(difflib.unified_diff(ta, tb, lineterm="", n=1))
    added = sum(1 for d in diff if d.startswith("+") and not d.startswith("+++"))
    removed = sum(1 for d in diff if d.startswith("-") and not d.startswith("---"))
    return {
        "verb": "compare-text",
        "inputs": [str(resolve(a.file)), str(resolve(a.other))],
        "outputs": [],
        "identical": added == 0 and removed == 0,
        "added": added,
        "removed": removed,
        "diff": diff[:200],
    }


def cmd_compare_visual(a):
    import pymupdf
    from PIL import ImageChops

    src_a, src_b = resolve(a.file), resolve(a.other)
    outdir = out_dir(a, src_a, "diff")
    da, db = pymupdf.open(src_a), pymupdf.open(src_b)
    outputs, changed = [], []
    for i in range(min(da.page_count, db.page_count)):
        ia = da[i].get_pixmap(dpi=a.dpi).pil_image().convert("RGB")
        ib = db[i].get_pixmap(dpi=a.dpi).pil_image().convert("RGB")
        if ia.size != ib.size:
            ib = ib.resize(ia.size)
        diff = ImageChops.difference(ia, ib)
        bbox = diff.getbbox()
        ratio = sum(diff.convert("L").point(lambda x: 1 if x else 0).getdata()) / (
            ia.width * ia.height
        )
        p = outdir / f"diff_p{i + 1}.png"
        diff.save(p)
        outputs.append(str(p))
        changed.append({"page": i + 1, "changed_ratio": round(ratio, 4), "bbox": bbox})
    da.close()
    db.close()
    return {
        "verb": "compare-visual",
        "inputs": [str(src_a), str(src_b)],
        "outputs": outputs,
        "pages": changed,
    }


# --------------------------------------------------------------------------- #
# Misc
# --------------------------------------------------------------------------- #
def cmd_repair(a):
    import shutil
    import subprocess

    src = resolve(a.file)
    out = ensure_parent(a.output or default_out(src, "repaired"))
    qpdf = shutil.which("qpdf")
    if not qpdf:
        raise PdfGoatError("repair requires qpdf on PATH")
    proc = subprocess.run(
        [qpdf, "--replace-input" if False else str(src), out],
        capture_output=True,
        text=True,
        check=False,
    )
    # qpdf returns 3 when it repairs a file with warnings; accept 0 or 3.
    if proc.returncode not in (0, 3) or not Path(out).exists():
        raise PdfGoatError(f"qpdf repair failed: {proc.stderr.strip()[:200]}")
    return {
        "verb": "repair",
        "inputs": [str(src)],
        "outputs": [out],
        "warnings": bool(proc.returncode == 3),
    }


def cmd_attach(a):
    import pymupdf

    src = resolve(a.file)
    att = resolve(a.attachment)
    out = ensure_parent(a.output or default_out(src, "attached"))
    doc = pymupdf.open(src)
    doc.embfile_add(att.name, att.read_bytes(), filename=att.name)
    _save_pdf(doc, out)
    doc.close()
    return {
        "verb": "attach",
        "inputs": [str(src)],
        "outputs": [out],
        "attached": att.name,
    }


def cmd_detach(a):
    import pymupdf

    src = resolve(a.file)
    names = list(dict.fromkeys(a.names or []))
    with pymupdf.open(src) as doc:
        catalog_names = doc.embfile_names()
        annotation_names = [name for _, _, name in _file_attachment_annotations(doc)]
        available = list(dict.fromkeys([*catalog_names, *annotation_names]))
        targets = available if a.remove_all else names
        missing = [name for name in targets if name not in available]
        if missing:
            raise PdfGoatError(f"attachment not found: {', '.join(missing)}")

        target_set = set(targets)
        removed_count = 0
        for name in catalog_names:
            if name in target_set:
                doc.embfile_del(name)
                removed_count += 1
        for page_index in range(doc.page_count):
            page = doc[page_index]
            annotation = page.first_annot
            while annotation:
                next_annotation = annotation.next
                if annotation.type[0] == pymupdf.PDF_ANNOT_FILE_ATTACHMENT:
                    name = annotation.file_info.get("filename") or ""
                    if name in target_set:
                        page.delete_annot(annotation)
                        removed_count += 1
                annotation = next_annotation

        out = ensure_parent(a.output or default_out(src, "detached"))
        doc.save(out, garbage=3, deflate=True)
    return {
        "verb": "detach",
        "inputs": [str(src)],
        "outputs": [out],
        "removed": targets,
        "count": removed_count,
    }


def cmd_search(a):
    import pymupdf

    src = resolve(a.file)
    limit = a.limit
    if limit is not None and limit < 1:
        raise PdfGoatError("--limit must be 1 or more")
    hits = []
    truncated = False
    with pymupdf.open(src) as doc:
        indices = list(
            parse_pages(a.pages, doc.page_count) if a.pages else range(doc.page_count)
        )
        for position, index in enumerate(indices):
            for rect in doc[index].search_for(a.query):
                hits.append({"page": index + 1, "rect": [round(v, 1) for v in rect]})
            if limit and len(hits) >= limit:
                truncated = len(hits) > limit or position + 1 < len(indices)
                del hits[limit:]
                break
    return {
        "verb": "search",
        "inputs": [str(src)],
        "outputs": [],
        "query": a.query,
        "count": len(hits),
        "hits": hits,
        "truncated": truncated,
    }


def cmd_overlay(a):
    import pymupdf

    src = resolve(a.file)
    stamp = resolve(a.stamp)
    out = ensure_parent(a.output or default_out(src, "overlay"))
    doc = pymupdf.open(src)
    sdoc = pymupdf.open(stamp)
    for i in range(doc.page_count):
        doc[i].show_pdf_page(
            doc[i].rect, sdoc, min(i, sdoc.page_count - 1), overlay=True
        )
    _save_pdf(doc, out)
    doc.close()
    return {"verb": "overlay", "inputs": [str(src), str(stamp)], "outputs": [out]}


def cmd_count(a):
    import pymupdf

    src = resolve(a.file)
    doc = pymupdf.open(src)
    words = chars = 0
    for index in range(doc.page_count):
        text, page_words = _page_text_and_words(doc[index])
        words += len(page_words)
        chars += len(text)
    pages = doc.page_count
    doc.close()
    return {
        "verb": "count",
        "inputs": [str(src)],
        "outputs": [],
        "pages": pages,
        "words": words,
        "chars": chars,
    }


# --------------------------------------------------------------------------- #
# Human formatting
# --------------------------------------------------------------------------- #
def render_human(result):
    verb = result.get("verb", "")
    if verb == "capabilities":
        print(f"{result['command_count']} commands")
        print(f"families  {', '.join(result['families'])}")
        print(f"commands  {', '.join(result['commands'])}")
        requested = result.get("requested_command")
        if requested:
            commands = sorted(result["schemas"][requested]["commands"])
            print(
                f"{requested}  {', '.join(commands) if commands else 'single command'}"
            )
        return
    if verb == "inspect":
        pages = result["pages"]
        end_page = result["start_page"] + len(pages) - 1
        print(
            f"{result['inputs'][0]}: pages {result['start_page']}-{end_page} "
            f"of {result['total_pages']}"
        )
        for page in pages:
            print(
                f"  {page['page']}: {page['width_pt']}x{page['height_pt']}pt, "
                f"{page['word_count']} words, {page['image_count']} images, "
                f"{page['link_count']} links"
            )
        if result["next_page"]:
            print(f"next page  {result['next_page']}")
        return
    if verb == "preflight":
        print(
            f"risk {result['risk']}  pages {result.get('pages', 'unknown')}  "
            f"attachments {result.get('attachments', 'unknown')}"
        )
        for finding in result["findings"]:
            count = f" ({finding['count']})" if "count" in finding else ""
            print(f"  {finding['severity']}: {finding['message']}{count}")
        return
    if verb == "info":
        print(f"file        {result['inputs'][0]}")
        print(f"pages       {result['pages']}")
        print(f"size        {human_size(result['file_size_bytes'])}")
        print(f"encrypted   {result['encrypted']}")
        print(
            f"has_forms   {result['has_forms']} ({result['form_field_count']} fields)"
        )
        print(f"has_text    {result['has_text']}")
        dims = ", ".join(
            f"{s['width']}x{s['height']}pt" for s in result["page_sizes_pt"]
        )
        print(f"page_sizes  {dims}")
        if result["metadata"]:
            print("metadata")
            for key, value in result["metadata"].items():
                print(f"  {key:<12}{value}")
        return
    if verb == "text":
        if "page_count" in result:
            print(
                f"# {result['char_count']} chars across {result['page_count']} pages"
                f" -> {result['outputs'][0]}"
            )
            return
        print(f"# {result['char_count']} chars across {len(result['pages'])} pages")
        print()
        for page in result["pages"]:
            print(f"----- page {page['page']} -----")
            print(page["text"].rstrip())
        return
    if verb == "form-list":
        print(f"{result['field_count']} form field(s):")
        for field in result["fields"]:
            print(f"  {field['name']}  [{field['type']}]  = {field['value']}")
        return
    if verb == "jobs":
        for job in result["jobs"]:
            outputs = f" -> {len(job['outputs'])} out" if job["outputs"] else ""
            print(
                f"#{job['id']:<4} {job['ts']}  {job['verb']:<11} {job['status']:<7}{outputs}"
            )
        return
    if verb == "compress":
        print(
            f"{human_size(result['original_bytes'])} -> {human_size(result['compressed_bytes'])}"
            f"  ({result['ratio']}x, saved {human_size(result['saved_bytes'])})"
        )
        print(f"out  {result['outputs'][0]}")
        return
    if verb == "convert-from-office":
        print(f"out  {result['outputs'][0]}")
        for warning in result["warnings"]:
            print(f"warn {warning}")
        print(f"engine={result['engine']}  pages={result['pages']}")
        return
    line = []
    for key, value in result.items():
        if key in ("inputs", "verb"):
            continue
        if key == "outputs":
            for output in value:
                print(f"out  {output}")
        else:
            line.append(f"{key}={value}")
    if line:
        print("  ".join(line))


# --------------------------------------------------------------------------- #
# Command group builders
# --------------------------------------------------------------------------- #
def _ns(sub, name, help):
    return sub.add_parser(name, help=help).add_subparsers(
        dest=f"{name}_cmd", required=True
    )


def _add_annotate(sub):
    ns = _ns(sub, "annotate", "annotations")
    for kind in ("highlight", "underline", "strikeout"):
        p = ns.add_parser(kind, help=f"{kind} text matching a string")
        p.add_argument("file")
        p.add_argument("--find", required=True)
        p.add_argument("--pages")
        p.add_argument("--color")
        p.add_argument("-o", "--output")
        p.set_defaults(func=cmd_annot_markup, kind=kind)
    p = ns.add_parser("note", help="sticky note")
    p.add_argument("file")
    p.add_argument("--page", type=int, default=1)
    p.add_argument("--at", required=True, help="x,y")
    p.add_argument("--text", required=True)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_annot_note)
    p = ns.add_parser("textbox", help="free-text box")
    p.add_argument("file")
    p.add_argument("--page", type=int, default=1)
    p.add_argument("--rect", required=True)
    p.add_argument("--text", required=True)
    p.add_argument("--size", type=float, default=12)
    p.add_argument("--color")
    p.add_argument("--fill")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_annot_textbox)
    for kind in ("rect", "circle"):
        p = ns.add_parser(kind, help=f"draw a {kind}")
        p.add_argument("file")
        p.add_argument("--page", type=int, default=1)
        p.add_argument("--rect", required=True)
        p.add_argument("--color")
        p.add_argument("--fill")
        p.add_argument("--width", type=float, default=1.5)
        p.add_argument("-o", "--output")
        p.set_defaults(func=cmd_annot_shape, kind=kind)
    for kind, arrow in (("line", False), ("arrow", True)):
        p = ns.add_parser(kind, help=f"draw a {kind}")
        p.add_argument("file")
        p.add_argument("--page", type=int, default=1)
        p.add_argument("--start", required=True, help="x,y")
        p.add_argument("--end", required=True, help="x,y")
        p.add_argument("--color")
        p.add_argument("--width", type=float, default=1.5)
        p.add_argument("-o", "--output")
        p.set_defaults(func=cmd_annot_line, arrow=arrow)
    p = ns.add_parser("ink", help="freehand ink stroke")
    p.add_argument("file")
    p.add_argument("--page", type=int, default=1)
    p.add_argument("--points", required=True, help="x,y;x,y;...")
    p.add_argument("--color")
    p.add_argument("--width", type=float, default=1.5)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_annot_ink)
    p = ns.add_parser("stamp", help="rubber stamp")
    p.add_argument("file")
    p.add_argument("--page", type=int, default=1)
    p.add_argument("--rect", required=True)
    p.add_argument(
        "--stamp", type=int, default=0, help="predefined stamp index from 0 to 13"
    )
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_annot_stamp)
    p = ns.add_parser("callout", help="callout box with arrow")
    p.add_argument("file")
    p.add_argument("--page", type=int, default=1)
    p.add_argument("--rect", required=True)
    p.add_argument("--target", required=True, help="arrow endpoint as x,y")
    p.add_argument("--text", required=True)
    p.add_argument("--size", type=float, default=11)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_annot_callout)
    p = ns.add_parser("area-highlight", help="add a transparent rectangular highlight")
    p.add_argument("file")
    p.add_argument("--page", type=int, default=1)
    p.add_argument("--rect", required=True)
    p.add_argument("--color")
    p.add_argument("--opacity", type=float, default=0.25, help="0 to 1")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_annot_area)
    p = ns.add_parser("polygon", help="draw a polygon")
    p.add_argument("file")
    p.add_argument("--page", type=int, default=1)
    p.add_argument("--points", required=True, help="x,y;x,y;x,y;...")
    p.add_argument("--color")
    p.add_argument("--fill")
    p.add_argument("--width", type=float, default=1.5)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_annot_polygon)
    p = ns.add_parser("list", help="list annotations")
    p.add_argument("file")
    p.set_defaults(func=cmd_annotations)
    p = ns.add_parser("flatten", help="flatten annotations into content")
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_annot_flatten)
    p = ns.add_parser("delete", help="delete annotations")
    p.add_argument("file")
    p.add_argument("--type", help="only this annotation type")
    p.add_argument("--pages")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_annot_delete)


def _add_security(sub):
    ns = _ns(sub, "security", "encryption, signatures, and sanitization")
    p = ns.add_parser("encrypt", help="encrypt with AES-256")
    p.add_argument("file")
    p.add_argument("--password", required=True, help="user password")
    p.add_argument("--owner")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_sec_encrypt)
    p = ns.add_parser("decrypt", help="remove encryption")
    p.add_argument("file")
    p.add_argument("--password", required=True)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_sec_decrypt)
    p = ns.add_parser("permissions", help="restrict printing, copying, or changes")
    p.add_argument("file")
    p.add_argument("--owner", required=True)
    p.add_argument("--user", default="")
    p.add_argument("--no-print", action="store_true")
    p.add_argument("--no-copy", action="store_true")
    p.add_argument("--no-modify", action="store_true")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_sec_permissions)
    p = ns.add_parser("sign", help="sign with a self-signed certificate")
    p.add_argument("file")
    p.add_argument("--name")
    p.add_argument("--reason")
    p.add_argument("--field")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_sec_sign)
    p = ns.add_parser("verify", help="verify signatures")
    p.add_argument("file")
    p.set_defaults(func=cmd_sec_verify)
    p = ns.add_parser(
        "sanitize",
        help="remove JavaScript, embedded or attached files, XMP metadata, and thumbnails",
    )
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_sec_sanitize)


def _add_meta(sub):
    ns = _ns(sub, "meta", "read and edit metadata")
    p = ns.add_parser("get", help="read metadata")
    p.add_argument("file")
    p.set_defaults(func=cmd_meta_get)
    p = ns.add_parser("set", help="set metadata key=value")
    p.add_argument("file")
    p.add_argument("--set", action="append", help="key=value (repeatable)")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_meta_set)
    p = ns.add_parser("strip", help="remove all metadata")
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_meta_strip)


def _add_pages(sub):
    ns = _ns(sub, "pages", "layout, numbering, imposition")
    p = ns.add_parser("blank", help="insert blank pages matching the adjacent page")
    p.add_argument("file")
    p.add_argument("--at", type=int, help="1-based insertion position; default: end")
    p.add_argument("--count", type=int, default=1)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_blank)
    p = ns.add_parser("duplicate", help="duplicate selected pages in place")
    p.add_argument("file")
    p.add_argument("--pages", required=True)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_duplicate)
    p = ns.add_parser("crop", help="set crop box")
    p.add_argument("file")
    p.add_argument("--box", required=True, help="x0,y0,x1,y1")
    p.add_argument("--pages")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_crop)
    p = ns.add_parser("scale", help="scale page size")
    p.add_argument("file")
    p.add_argument("--factor", type=float, required=True)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_scale)
    p = ns.add_parser("nup", help="n-up imposition (2 or 4)")
    p.add_argument("file")
    p.add_argument("--n", type=int, choices=[2, 4], required=True)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_nup)
    p = ns.add_parser("booklet", help="saddle-stitch booklet imposition")
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_booklet)
    for where in ("header", "footer"):
        p = ns.add_parser(where, help=f"add a {where}")
        p.add_argument("file")
        p.add_argument("--text", required=True, help="supports {page} {pages}")
        p.add_argument("--align", default="center", choices=["left", "center", "right"])
        p.add_argument("--size", type=float, default=10)
        p.add_argument("--color")
        p.add_argument("-o", "--output")
        p.set_defaults(func=cmd_pages_header, where=where)
    p = ns.add_parser("numbers", help="add page numbers")
    p.add_argument("file")
    p.add_argument("--format", default="{page}", help="e.g. 'Page {page} of {pages}'")
    p.add_argument("--start", type=int, default=1)
    p.add_argument("--align", default="center", choices=["left", "center", "right"])
    p.add_argument("--size", type=float, default=10)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_numbers)
    p = ns.add_parser("bates", help="Bates numbering")
    p.add_argument("file")
    p.add_argument("--prefix", default="")
    p.add_argument("--start", type=int, default=1)
    p.add_argument("--digits", type=int, default=6)
    p.add_argument("--size", type=float, default=9)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_bates)
    p = ns.add_parser("boxes", help="set the media, crop, trim, or bleed box")
    p.add_argument("file")
    p.add_argument("--box", required=True, choices=["media", "crop", "trim", "bleed"])
    p.add_argument("--rect", required=True)
    p.add_argument("--pages")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_boxes)
    p = ns.add_parser("insert", help="insert pages from another PDF")
    p.add_argument("file")
    p.add_argument("--source", required=True)
    p.add_argument("--at", type=int, help="1-based position (default: end)")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_insert)
    p = ns.add_parser("replace", help="replace pages with another PDF")
    p.add_argument("file")
    p.add_argument("--source", required=True)
    p.add_argument("--pages", required=True)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_replace)
    p = ns.add_parser(
        "flatten",
        help="flatten annotations and form appearances; flatten transparency when Ghostscript succeeds",
    )
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_pages_flatten)


def _add_get(sub):
    ns = _ns(sub, "get", "extract assets")
    p = ns.add_parser("images", help="extract embedded images")
    p.add_argument("file")
    p.add_argument("-o", "--outdir")
    p.set_defaults(func=cmd_get_images)
    p = ns.add_parser("fonts", help="list fonts")
    p.add_argument("file")
    p.set_defaults(func=cmd_get_fonts)
    p = ns.add_parser("text-blocks", help="extract bounded text with page geometry")
    p.add_argument("file")
    p.add_argument("--pages", help="default: first 25 pages")
    p.add_argument("--max-blocks", type=int, default=200)
    p.add_argument("--start-block", type=int, default=0)
    p.set_defaults(func=cmd_get_text_blocks)
    p = ns.add_parser("attachments", help="extract embedded files")
    p.add_argument("file")
    p.add_argument("-o", "--outdir")
    p.set_defaults(func=cmd_get_attachments)
    p = ns.add_parser("bookmarks", help="list document outline entries")
    p.add_argument("file")
    p.set_defaults(func=cmd_get_bookmarks)
    p = ns.add_parser("links", help="list links")
    p.add_argument("file")
    p.set_defaults(func=cmd_get_links)


def _add_nav(sub):
    ns = _ns(sub, "bookmarks", "edit the document outline")
    p = ns.add_parser("set", help="set outline from JSON [{level,title,page}]")
    p.add_argument("file")
    p.add_argument("--data", required=True)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_bookmarks_set)
    p = ns.add_parser("clear", help="remove the document outline")
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_bookmarks_clear)
    ns2 = _ns(sub, "links", "edit links")
    p = ns2.add_parser("add", help="add a link")
    p.add_argument("file")
    p.add_argument("--page", type=int, default=1)
    p.add_argument("--rect", required=True)
    p.add_argument("--uri")
    p.add_argument("--goto", type=int, help="target page (1-based)")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_links_add)
    p = ns2.add_parser("remove", help="remove links")
    p.add_argument("file")
    p.add_argument("--pages")
    p.add_argument("--external-only", action="store_true")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_links_remove)


def _add_convert(sub):
    ns = _ns(sub, "convert", "conversion and OCR")
    p = ns.add_parser("ocr", help="add a searchable text layer (ocrmypdf)")
    p.add_argument("file")
    p.add_argument("--force", action="store_true", help="re-OCR even if text exists")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_convert_ocr)
    p = ns.add_parser("html", help="convert PDF pages to HTML")
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_convert_html)
    p = ns.add_parser("tables", help="extract tables to CSV (pdfplumber)")
    p.add_argument("file")
    p.add_argument("-o", "--outdir")
    p.set_defaults(func=cmd_convert_tables)
    p = ns.add_parser(
        "pdfa", help="create an unvalidated PDF/A-2b candidate with Ghostscript"
    )
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_convert_pdfa)
    p = ns.add_parser("docx", help="convert PDF to Word (.docx) with pdf2docx")
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_convert_docx)
    p = ns.add_parser("xlsx", help="convert extracted PDF tables to Excel (.xlsx)")
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_convert_xlsx)
    p = ns.add_parser("pptx", help="render each PDF page as one PowerPoint slide")
    p.add_argument("file")
    p.add_argument("--dpi", type=int, default=150)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_convert_pptx)
    p = ns.add_parser(
        "from-office",
        help="convert a .docx, .xlsx, or .pptx file to PDF with office2pdf",
    )
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_convert_from_office)
    p = ns.add_parser("audio", help="export extracted text as AIFF with macOS say")
    p.add_argument("file")
    p.add_argument("--voice")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_convert_audio)


def _add_optimize(sub):
    ns = _ns(sub, "optimize", "reduce file size")
    p = ns.add_parser("reduce", help="reduce file size (ghostscript presets)")
    p.add_argument("file")
    p.add_argument(
        "--preset", default="ebook", choices=["screen", "ebook", "printer", "prepress"]
    )
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_optimize_reduce)


def _add_edit(sub):
    ns = _ns(sub, "edit", "edit page content")
    p = ns.add_parser(
        "text",
        help="find and replace simple text runs; no reflow or embedded-font matching",
    )
    p.add_argument("file")
    p.add_argument("--find", required=True)
    p.add_argument("--replace", required=True)
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_edit_text)


def _add_access(sub):
    ns = _ns(sub, "accessibility", "check and set accessibility metadata")
    p = ns.add_parser(
        "check", help="report basic tag, title, language, and image-alt checks"
    )
    p.add_argument("file")
    p.set_defaults(func=cmd_access_check)
    p = ns.add_parser("set", help="set title and language; set the Marked flag")
    p.add_argument("file")
    p.add_argument("--title")
    p.add_argument("--lang", default="en")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_access_set)


def _add_compare(sub):
    ns = _ns(sub, "compare", "compare two PDFs")
    p = ns.add_parser("text", help="compare extracted text")
    p.add_argument("file")
    p.add_argument("other")
    p.set_defaults(func=cmd_compare_text)
    p = ns.add_parser("visual", help="write one visual-difference PNG per page")
    p.add_argument("file")
    p.add_argument("other")
    p.add_argument("--dpi", type=int, default=100)
    p.add_argument("-o", "--outdir")
    p.set_defaults(func=cmd_compare_visual)


def _add_misc(sub):
    p = sub.add_parser("repair", help="repair a damaged PDF (qpdf)")
    p.add_argument("file")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_repair)
    p = sub.add_parser("attach", help="embed a file attachment")
    p.add_argument("file")
    p.add_argument("attachment")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_attach)
    p = sub.add_parser("detach", help="remove embedded file attachments")
    p.add_argument("file")
    target = p.add_mutually_exclusive_group(required=True)
    target.add_argument("--name", dest="names", action="append")
    target.add_argument("--all", dest="remove_all", action="store_true")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_detach)
    p = sub.add_parser("search", help="find text and return PDF-point rectangles")
    p.add_argument("file")
    p.add_argument("query")
    p.add_argument(
        "--first",
        action="store_const",
        const=1,
        dest="limit",
        help="stop at the first hit",
    )
    p.add_argument("--limit", type=int, help="stop after this many hits")
    p.add_argument("--pages", help="default: all pages")
    p.set_defaults(func=cmd_search)
    p = sub.add_parser("overlay", help="stamp one PDF over another")
    p.add_argument("file")
    p.add_argument("stamp")
    p.add_argument("-o", "--output")
    p.set_defaults(func=cmd_overlay)
    p = sub.add_parser("count", help="count pages, words, and characters")
    p.add_argument("file")
    p.set_defaults(func=cmd_count)


# --------------------------------------------------------------------------- #
# CLI
# --------------------------------------------------------------------------- #
def cmd_transcript_read(a):
    from .transcript import parse_transcript

    src = resolve(a.file)
    output_path = None
    if a.output:
        output_path = Path(a.output).expanduser().resolve()
        if output_path == src or output_path.exists() and output_path.samefile(src):
            raise PdfGoatError("output must differ from input")
    parsed = parse_transcript(src, a.conferred)
    parsed.pop("layout", None)
    result = {
        "verb": "transcript-read",
        "inputs": [str(src)],
        "outputs": [],
        **parsed,
    }
    if output_path is not None:
        out = ensure_parent(output_path)
        Path(out).write_text(json.dumps(parsed, indent=2, default=str))
        result["outputs"] = [out]
    return result


def cmd_transcript_resolve(a):
    from .transcript import discover_transcripts

    rows = discover_transcripts(a.root, a.glob)
    return {
        "verb": "transcript-resolve",
        "inputs": [str(Path(root).expanduser().resolve()) for root in a.root],
        "outputs": [],
        "roots": [str(Path(root).expanduser().resolve()) for root in a.root],
        "candidates": rows,
        "candidate_count": len(rows),
    }


def _add_transcript(sub):
    ns = _ns(sub, "transcript", "extract academic transcript data")
    p = ns.add_parser(
        "read", help="extract transcript identity, terms, courses, and freshness"
    )
    p.add_argument("file")
    p.add_argument(
        "--conferred",
        help="conferral date to compare with the printed issue date (YYYY-MM-DD)",
    )
    p.add_argument("-o", "--output", help="write structured JSON")
    p.set_defaults(func=cmd_transcript_read, output=None)
    p = ns.add_parser(
        "resolve",
        help="rank PDF candidates by printed issue date in the named files or directories",
    )
    p.add_argument(
        "--root", action="append", required=True, help="directory or PDF; not recursive"
    )
    p.add_argument(
        "--glob", action="append", dest="glob", help="filename pattern, repeatable"
    )
    p.set_defaults(func=cmd_transcript_resolve, glob=None)


def build_parser():
    p = PdfGoatArgumentParser(
        prog="pdf-goat", description="Local PDF editing and inspection tool"
    )
    # `capabilities` reports this parser instead of building a second copy.
    p.set_defaults(root_parser=p, ledger=True)
    p.add_argument(
        "--agent", action="store_true", help="write JSON output even on a TTY"
    )
    sub = p.add_subparsers(dest="cmd", required=True)
    _add_transcript(sub)

    s = sub.add_parser("capabilities", help="discover command schemas for agents")
    s.add_argument("family", nargs="?", help="top-level command family")
    s.set_defaults(func=cmd_capabilities, ledger=False)

    s = sub.add_parser("inspect", help="list page sizes and content counts")
    s.add_argument("file")
    s.add_argument("--start-page", type=int, default=1)
    s.add_argument("--limit", type=int, default=_DEFAULT_PAGE_WINDOW)
    s.set_defaults(func=cmd_inspect)

    s = sub.add_parser(
        "preflight",
        help="inspect active content, links, forms, attachments, and basic accessibility signals",
    )
    s.add_argument("file")
    s.set_defaults(func=cmd_preflight)

    s = sub.add_parser("info", help="show document metadata and page summary")
    s.add_argument("file")
    s.set_defaults(func=cmd_info)

    s = sub.add_parser("merge", help="concatenate PDFs")
    s.add_argument("files", nargs="+")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_merge)

    s = sub.add_parser("split", help="split a PDF into chunks")
    s.add_argument("file")
    s.add_argument("--every", type=int, default=1, help="pages per chunk")
    s.add_argument("-o", "--outdir")
    s.set_defaults(func=cmd_split)

    s = sub.add_parser("extract", help="extract pages, e.g. 2-5,9")
    s.add_argument("file")
    s.add_argument("--pages", required=True)
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_extract)

    s = sub.add_parser("delete", help="delete pages")
    s.add_argument("file")
    s.add_argument("--pages", required=True)
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_delete)

    s = sub.add_parser("reorder", help="reorder pages, e.g. --order 3,1,2")
    s.add_argument("file")
    s.add_argument("--order", required=True)
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_reorder)

    s = sub.add_parser("rotate", help="rotate pages by a multiple of 90")
    s.add_argument("file")
    s.add_argument("--pages", help="default: all pages")
    s.add_argument("--deg", type=int, required=True)
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_rotate)

    s = sub.add_parser("render", help="rasterize pages to images")
    s.add_argument("file")
    s.add_argument("--pages", help="default: all pages")
    s.add_argument("--dpi", type=int, default=150)
    s.add_argument("--format", default="png", choices=["png", "jpg", "ppm"])
    s.add_argument("--clip", help="render only x0,y0,x1,y1 in PDF points")
    s.add_argument("-o", "--outdir")
    s.set_defaults(func=cmd_render)

    s = sub.add_parser("from-images", help="build a PDF from images")
    s.add_argument("images", nargs="+")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_from_images)

    s = sub.add_parser("redact", help="redact words matching a regular expression")
    s.add_argument("file")
    s.add_argument("--find", required=True, help="regex matched per word")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_redact)

    s = sub.add_parser("watermark", help="stamp a diagonal text watermark")
    s.add_argument("file")
    s.add_argument("--text", default="DRAFT")
    s.add_argument("--size", type=float, default=72)
    s.add_argument("--opacity", type=float, default=0.15)
    s.add_argument("--angle", type=float, default=45)
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_watermark)

    s = sub.add_parser(
        "compress",
        help="recompress and linearize; keep the input bytes if the result would be larger",
    )
    s.add_argument("file")
    s.add_argument(
        "--level",
        default="/ebook",
        choices=["/screen", "/ebook", "/printer", "/prepress"],
    )
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_compress)

    s = sub.add_parser("text", help="extract text")
    s.add_argument("file")
    s.add_argument("-o", "--output", help="write text to a file")
    s.add_argument(
        "--layout", action="store_true", help="preserve positioned columns and lines"
    )
    s.set_defaults(func=cmd_text)

    s = sub.add_parser("from-html", help="render an HTML file to PDF (weasyprint)")
    s.add_argument("file")
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_from_html)

    s = sub.add_parser("from-md", help="render Markdown to a styled PDF")
    s.add_argument("file")
    s.add_argument(
        "--css", help="path to a CSS file (overrides the default stylesheet)"
    )
    s.add_argument("-o", "--output")
    s.set_defaults(func=cmd_from_md)

    s = sub.add_parser("form", help="form fields")
    fsub = s.add_subparsers(dest="form_cmd", required=True)
    fl = fsub.add_parser("list", help="list form fields")
    fl.add_argument("file")
    fl.set_defaults(func=cmd_form_fields)
    ff = fsub.add_parser("fill", help="fill a form from JSON")
    ff.add_argument("file")
    ff.add_argument("--data", required=True, help="JSON file of field:value")
    ff.add_argument("--flatten", action="store_true")
    ff.add_argument("-o", "--output")
    ff.set_defaults(func=cmd_form_fill)
    fct = fsub.add_parser("create-text", help="add a text field")
    fct.add_argument("file")
    fct.add_argument("--name", required=True)
    fct.add_argument("--page", type=int, default=1)
    fct.add_argument("--rect", required=True, help="x0,y0,x1,y1")
    fct.add_argument("-o", "--output")
    fct.set_defaults(func=cmd_form_create_text)
    fcc = fsub.add_parser("create-checkbox", help="add a checkbox field")
    fcc.add_argument("file")
    fcc.add_argument("--name", required=True)
    fcc.add_argument("--page", type=int, default=1)
    fcc.add_argument("--rect", required=True, help="x0,y0,x1,y1")
    fcc.add_argument("-o", "--output")
    fcc.set_defaults(func=cmd_form_create_checkbox)
    fe = fsub.add_parser("export", help="export field data")
    fe.add_argument("file")
    fe.add_argument("--format", default="json", choices=["json", "xfdf", "fdf"])
    fe.add_argument("-o", "--output")
    fe.set_defaults(func=cmd_form_export)
    fi = fsub.add_parser("import", help="import JSON or XFDF field data")
    fi.add_argument("file")
    fi.add_argument("--data", required=True)
    fi.add_argument("--flatten", action="store_true")
    fi.add_argument("-o", "--output")
    fi.set_defaults(func=cmd_form_import)

    _add_annotate(sub)
    _add_security(sub)
    _add_meta(sub)
    _add_pages(sub)
    _add_get(sub)
    _add_nav(sub)
    _add_convert(sub)
    _add_optimize(sub)
    _add_edit(sub)
    _add_access(sub)
    _add_compare(sub)
    _add_misc(sub)

    s = sub.add_parser("jobs", help="show the job ledger")
    s.add_argument("--limit", type=int, default=20)
    s.set_defaults(func=cmd_jobs, ledger=False)

    return p


_LEDGER_COUNT_NAMES = {
    "candidates": "candidate_count",
    "courses": "course_count",
    "terms": "term_count",
}


def _ledger_detail(result):
    detail = {}
    for key, value in result.items():
        if key in ("inputs", "outputs", "verb"):
            continue
        if key == "freshness" and isinstance(value, dict):
            detail["freshness_verdict"] = value.get("verdict")
            continue
        if key == "parse_quality" and isinstance(value, dict):
            for nested_key in (
                "confidence",
                "matched_line_count",
                "unparsed_line_count",
                "course_count",
                "term_count",
            ):
                if nested_key in value:
                    detail[nested_key] = value[nested_key]
            continue
        match value:
            case None | bool() | int() | float():
                detail[key] = value
            case str():
                if key == "risk" and value in {"low", "medium", "high"}:
                    detail[key] = value
                else:
                    detail[f"{key}_char_count"] = len(value)
            case list() | tuple() | set() | dict():
                count_key = _LEDGER_COUNT_NAMES.get(key, f"{key}_count")
                if key == "pages" and "total_pages" not in result:
                    count_key = "page_count"
                detail[count_key] = len(value)
    return detail


def main(argv=None):
    parser = build_parser()
    tokens = list(sys.argv[1:] if argv is None else argv)
    json_mode = "--agent" in tokens or not sys.stdout.isatty()
    verb = "pdf-goat"
    ledger = False
    parsed_arguments = False
    start = time.time()
    try:
        args = parser.parse_args(tokens)
        parsed_arguments = True
        verb = args.cmd
        ledger = args.ledger
        result = args.func(args)
        status = "success"
        message = None
    except PdfGoatError as error:
        result = {"verb": verb, "inputs": [], "outputs": [], "error": str(error)}
        status, message = "error", str(error)
    except Exception as error:  # noqa: BLE001
        result = {
            "verb": verb,
            "inputs": [],
            "outputs": [],
            "error": f"{type(error).__name__}: {error}",
        }
        status, message = "error", str(error)
    duration = int((time.time() - start) * 1000)

    if ledger and parsed_arguments:
        detail = _ledger_detail(result)
        record_job(
            result.get("verb", verb),
            status,
            result.get("inputs", []),
            result.get("outputs", []),
            detail,
            message,
            duration,
        )

    result["ok"] = status == "success"
    if json_mode:
        print(json.dumps(result, indent=2, default=str))
    else:
        if status == "success":
            render_human(result)
        else:
            print(f"error: {result['error']}", file=sys.stderr)
    return 0 if status == "success" else 1


if __name__ == "__main__":
    sys.exit(main())
